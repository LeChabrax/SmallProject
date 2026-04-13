"""
Générateur de deck PowerPoint — Benchmark Abonnement Impulse Nutrition
python-pptx + Pillow | Couleurs Impulse : #1B2B4B (marine) + #F05A28 (orange)
"""

import json
import os
from pathlib import Path
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree

# ─── Constantes design ────────────────────────────────────────────────
MARINE   = RGBColor(0x1B, 0x2B, 0x4B)
ORANGE   = RGBColor(0xF0, 0x5A, 0x28)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xF2, 0xF4, 0xF7)
DGRAY    = RGBColor(0x4A, 0x4A, 0x4A)
GREEN    = RGBColor(0x2D, 0xA8, 0x6A)
RED_CLR  = RGBColor(0xE5, 0x3E, 0x3E)
AMBER    = RGBColor(0xF5, 0xA6, 0x23)

SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)

BASE_DIR = Path(".")
LOGO_PATH     = BASE_DIR / "assets/logo_impulse.jpeg"
SCREENSHOTS   = BASE_DIR / "assets/screenshots/benchmark"
DATA_PATH     = BASE_DIR / "benchmark/data.json"
RENTA_PATH    = BASE_DIR / "benchmark/rentabilite_output.json"
OUTPUT_PATH   = BASE_DIR / "Benchmark_Abonnement_Impulse_Avril2026.pptx"

FONT_TITLE   = "Calibri"
FONT_BODY    = "Calibri"

# ─── Helpers ──────────────────────────────────────────────────────────

def add_bg(slide, color=MARINE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=12, font_color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, italic=False, bg_color=None,
             font_name=None, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    else:
        txBox.fill.background()
    txBox.line.fill.background()

    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name or FONT_BODY
    return txBox


def add_multiline_text(slide, lines, left, top, width, height,
                        font_size=11, font_color=WHITE, line_spacing=Pt(4)):
    """lines = list of (text, bold, color) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.fill.background()
    txBox.line.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, bold, color = line_data, False, font_color
        else:
            text = line_data.get("text", "")
            bold = line_data.get("bold", False)
            color = line_data.get("color", font_color)

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)

        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT_BODY


def score_badge(slide, score, left, top, size=Inches(0.65)):
    """Affiche un badge score coloré /5."""
    if score >= 4.0:
        color = GREEN
    elif score >= 3.0:
        color = AMBER
    else:
        color = RED_CLR

    add_rect(slide, left, top, size, size * 0.65, color)
    add_text(slide, f"{score:.1f}/5", left, top, size, size * 0.65,
             font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def crop_screenshot(src_path, dst_path, max_height_px=900):
    """Recadre le screenshot pour garder seulement le haut de la page."""
    try:
        img = Image.open(src_path)
        w, h = img.size
        if h > max_height_px:
            img = img.crop((0, 0, w, max_height_px))
        img.save(dst_path, "PNG", optimize=True)
        return dst_path
    except Exception as e:
        print(f"  ⚠️  Crop error {src_path}: {e}")
        return src_path


def add_logo(slide, left=Inches(0.15), top=Inches(0.1), height=Inches(0.55)):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), left, top, height=height)


# ─── Slide 1 : Cover ──────────────────────────────────────────────────

def slide_cover(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    add_bg(slide, MARINE)

    # Bande orange gauche
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, ORANGE)

    # Logo en haut à droite
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(10.5), Inches(0.3), height=Inches(0.8))

    # Titre principal
    add_text(slide, "Benchmark", Inches(0.7), Inches(1.5), Inches(9), Inches(1.0),
             font_size=40, font_color=WHITE, bold=True, font_name=FONT_TITLE)
    add_text(slide, "Modèles d'Abonnement", Inches(0.7), Inches(2.35), Inches(11), Inches(1.0),
             font_size=40, font_color=ORANGE, bold=True, font_name=FONT_TITLE)
    add_text(slide, "Nutrition Sportive", Inches(0.7), Inches(3.2), Inches(9), Inches(0.8),
             font_size=36, font_color=WHITE, bold=False, font_name=FONT_TITLE)

    # Ligne séparatrice
    add_rect(slide, Inches(0.7), Inches(4.2), Inches(5), Inches(0.04), ORANGE)

    # Sous-titre
    add_text(slide, "Analyse stratégique & recommandation pour le lancement Impulse",
             Inches(0.7), Inches(4.45), Inches(11), Inches(0.5),
             font_size=14, font_color=LGRAY, font_name=FONT_TITLE)

    # Footer
    add_rect(slide, 0, Inches(6.9), SLIDE_W, Inches(0.6), RGBColor(0x0F, 0x1A, 0x2E))
    add_text(slide, "Avril 2026  |  Confidentiel — Usage interne comité projet",
             Inches(0.7), Inches(7.0), Inches(10), Inches(0.4),
             font_size=10, font_color=LGRAY, align=PP_ALIGN.LEFT)

    return slide


# ─── Slide 2 : Contexte & Objectif ────────────────────────────────────

def slide_contexte(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    # Header marine
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), MARINE)
    add_logo(slide, left=Inches(0.2), top=Inches(0.28), height=Inches(0.65))
    add_text(slide, "Contexte & Objectif", Inches(1.6), Inches(0.35), Inches(10), Inches(0.6),
             font_size=24, font_color=WHITE, bold=True)

    # 3 colonnes "questions clés"
    questions = [
        ("🔍", "Qu'est-ce qui fait un bon modèle\nd'abonnement nutrition sportive ?",
         "Benchmark de 8 marques FR/international\nsur 4 axes : offre, périmètre, UX, flexibilité"),
        ("💰", "À partir de quel panier la\nlivraison offerte est-elle rentable ?",
         "Modélisation financière avec 3 scénarios\nprix × marge × coût logistique"),
        ("🚀", "Quel modèle concret\nrecommander à Impulse ?",
         "Produits éligibles, remise, fréquences,\nwording et positionnement UX"),
    ]

    for i, (emoji, q, a) in enumerate(questions):
        x = Inches(0.4 + i * 4.3)
        # Carte bleue
        add_rect(slide, x, Inches(1.4), Inches(4.0), Inches(4.5), LGRAY)
        add_rect(slide, x, Inches(1.4), Inches(4.0), Inches(0.08), ORANGE)
        add_text(slide, emoji, x + Inches(0.15), Inches(1.55), Inches(0.7), Inches(0.7), font_size=28, font_color=MARINE)
        add_text(slide, q, x + Inches(0.15), Inches(2.3), Inches(3.7), Inches(1.1),
                 font_size=13, font_color=MARINE, bold=True)
        add_text(slide, a, x + Inches(0.15), Inches(3.5), Inches(3.7), Inches(2.0),
                 font_size=11, font_color=DGRAY)

    # Note en bas
    add_rect(slide, 0, Inches(6.1), SLIDE_W, Inches(0.08), ORANGE)
    add_text(slide,
             "Périmètre : 8 marques analysées  •  Nutri&Co  •  Nutrimuscle  •  Novoma  •  Nutripure  •  Cuure  •  Decathlon  •  Aqeelab  •  MyProtein",
             Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.6),
             font_size=10, font_color=DGRAY, align=PP_ALIGN.CENTER)


# ─── Slide 3 : Méthodologie ───────────────────────────────────────────

def slide_methodologie(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), MARINE)
    add_logo(slide, left=Inches(0.2), top=Inches(0.28), height=Inches(0.65))
    add_text(slide, "Méthodologie & Périmètre d'analyse", Inches(1.6), Inches(0.35), Inches(10), Inches(0.6),
             font_size=24, font_color=WHITE, bold=True)

    axes = [
        ("1", "Modèle\ncommercial", "Remise %, fréquences,\nlimite d'engagement,\nlivraison, flexibilité"),
        ("2", "Périmètre\nproduit", "Gammes éligibles,\nexclusions, ticket\nmoyen estimé"),
        ("3", "UX &\nfiche produit", "Module abo, wording,\npage dédiée, homepage,\nnavigation"),
        ("4", "Évaluation\nsynthétique", "Points forts/faibles,\nscore /5 commercial\n+ /5 UX → score global"),
    ]

    for i, (num, title, detail) in enumerate(axes):
        x = Inches(0.35 + i * 3.25)
        add_rect(slide, x, Inches(1.5), Inches(3.0), Inches(3.5), LGRAY)
        add_rect(slide, x, Inches(1.5), Inches(0.5), Inches(3.5), MARINE)
        add_text(slide, num, x + Inches(0.1), Inches(1.55), Inches(0.4), Inches(0.7),
                 font_size=22, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(0.15), Inches(2.25), Inches(2.7), Inches(0.7),
                 font_size=13, font_color=MARINE, bold=True)
        add_text(slide, detail, x + Inches(0.15), Inches(2.95), Inches(2.7), Inches(1.8),
                 font_size=10, font_color=DGRAY)

    # Scoring grille
    add_text(slide, "Grille de scoring", Inches(0.4), Inches(5.3), Inches(6), Inches(0.4),
             font_size=13, font_color=MARINE, bold=True)
    scoring_lines = [
        "Score commercial /5 : Remise (2pts) + Flexibilité (1.5pts) + Fréquences (1pt) + Livraison (0.5pt)",
        "Score UX /5 : Module produit (1pt) + Wording (1pt) + Page dédiée (1pt) + Homepage (1pt) + Navigation (1pt)",
        "Score global = moyenne des deux scores",
    ]
    for i, line in enumerate(scoring_lines):
        add_text(slide, f"• {line}", Inches(0.4), Inches(5.75 + i * 0.3), Inches(12.5), Inches(0.3),
                 font_size=10, font_color=DGRAY)

    add_text(slide, "Sources : sites officiels, FAQ, pages abonnement  •  Date de collecte : Avril 2026",
             Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.35),
             font_size=9, font_color=RGBColor(0xAA, 0xAA, 0xAA), italic=True)


# ─── Slide marque ─────────────────────────────────────────────────────

def slide_brand(prs, brand_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    brand_name = brand_data["brand"]
    has_sub = brand_data["has_subscription"]
    score_g = brand_data["evaluation"].get("score_global", 0) or 0
    score_c = brand_data["evaluation"].get("score_commercial", 0) or 0
    score_u = brand_data["evaluation"].get("score_ux", 0) or 0

    # Header
    header_color = MARINE if has_sub else RGBColor(0x55, 0x55, 0x55)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.05), header_color)
    add_logo(slide, left=Inches(0.2), top=Inches(0.2), height=Inches(0.6))
    add_text(slide, brand_name, Inches(1.6), Inches(0.15), Inches(7), Inches(0.55),
             font_size=26, font_color=WHITE, bold=True)

    # Badge abonnement
    status_txt = "✅ Abonnement actif" if has_sub else "❌ Pas d'abonnement"
    status_col = GREEN if has_sub else RED_CLR
    add_rect(slide, Inches(9.0), Inches(0.2), Inches(2.8), Inches(0.45), status_col)
    add_text(slide, status_txt, Inches(9.0), Inches(0.2), Inches(2.8), Inches(0.45),
             font_size=11, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Scores
    add_text(slide, f"Commercial", Inches(11.95), Inches(0.05), Inches(1.2), Inches(0.3),
             font_size=8, font_color=LGRAY, align=PP_ALIGN.CENTER)
    add_text(slide, f"{score_c:.1f}/5", Inches(11.95), Inches(0.3), Inches(1.2), Inches(0.4),
             font_size=14, font_color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

    # ── Colonne gauche: screenshot ──
    left_w = Inches(6.5)
    screenshot_path = SCREENSHOTS / f"{brand_name.lower().replace('&', '').replace(' ', '').replace('é', 'e').replace('/', '').replace('myprotein', 'myprotein')}.png"
    # Map brand name to file name
    brand_file_map = {
        "Nutri&Co": "nutriandco",
        "Nutrimuscle": "nutrimuscle",
        "Novoma": "novoma",
        "Nutripure": "nutripure",
        "Cuure": "cuure",
        "Decathlon": "decathlon",
        "Aqeelab": "aqeelab",
        "MyProtein": "myprotein",
    }
    fname = brand_file_map.get(brand_name, brand_name.lower())
    screenshot_path = SCREENSHOTS / f"{fname}.png"
    cropped_path = SCREENSHOTS / f"{fname}_crop.png"

    if screenshot_path.exists():
        try:
            crop_screenshot(str(screenshot_path), str(cropped_path), max_height_px=800)
            slide.shapes.add_picture(str(cropped_path), Inches(0.15), Inches(1.15),
                                     width=left_w - Inches(0.2), height=Inches(5.7))
        except Exception as e:
            print(f"  ⚠️  Image error {fname}: {e}")
            add_rect(slide, Inches(0.15), Inches(1.15), left_w - Inches(0.2), Inches(5.7), LGRAY)
            add_text(slide, f"📸 {brand_name}\n[capture écran]",
                     Inches(1.5), Inches(3.5), Inches(4), Inches(1.2),
                     font_size=14, font_color=DGRAY, align=PP_ALIGN.CENTER)
    else:
        add_rect(slide, Inches(0.15), Inches(1.15), left_w - Inches(0.2), Inches(5.7), LGRAY)
        add_text(slide, "Screenshot non disponible", Inches(1.5), Inches(3.5), Inches(4), Inches(0.5),
                 font_size=11, font_color=DGRAY, align=PP_ALIGN.CENTER)

    # ── Colonne droite: données ──
    rx = left_w + Inches(0.15)
    rw = SLIDE_W - rx - Inches(0.15)
    ry = Inches(1.15)

    cm = brand_data.get("commercial_model", {})

    def fmt(v, default="—"):
        if v is None or v == []:
            return default
        if isinstance(v, list):
            return " / ".join(str(x) for x in v)
        if isinstance(v, bool):
            return "✅ Oui" if v else "❌ Non"
        return str(v)

    # Tableau données commerciales
    rows = []
    if has_sub:
        rows = [
            ("Remise", fmt(cm.get("discount_detail") or cm.get("discount_pct") and f"{cm['discount_pct']}%")),
            ("Fréquences", fmt(cm.get("frequencies"))),
            ("Engagement", fmt(cm.get("min_engagement"))),
            ("Livraison", "✅ Offerte" if cm.get("free_delivery") else f"❌ {cm.get('delivery_cost_below', 'payante')}"),
            ("Annulation", fmt(cm.get("cancellation_policy", "")[:50])),
        ]
    else:
        alt = brand_data.get("loyalty_alternative", {})
        rows = [
            ("Abonnement", "❌ Absent"),
            ("Alternative", fmt(alt.get("program_name"))),
            ("Fidélité", fmt(alt.get("description", "")[:80])),
            ("Hypothèse", fmt(brand_data.get("strategic_hypothesis", "")[:80])),
        ]

    # Tableau
    for i, (label, value) in enumerate(rows):
        row_y = ry + Inches(i * 0.55)
        bg = LGRAY if i % 2 == 0 else WHITE
        add_rect(slide, rx, row_y, rw, Inches(0.5), bg)
        add_text(slide, label, rx + Inches(0.1), row_y + Inches(0.05), Inches(1.3), Inches(0.45),
                 font_size=9, font_color=MARINE, bold=True)
        add_text(slide, value, rx + Inches(1.4), row_y + Inches(0.05), rw - Inches(1.5), Inches(0.45),
                 font_size=9, font_color=DGRAY)

    # Séparateur
    sep_y = ry + Inches(len(rows) * 0.55 + 0.1)
    add_rect(slide, rx, sep_y, rw, Inches(0.03), ORANGE)

    # Points forts
    strengths = brand_data["evaluation"].get("strengths", [])
    weaknesses = brand_data["evaluation"].get("weaknesses", [])
    bullet_y = sep_y + Inches(0.1)

    if strengths:
        add_text(slide, "Points forts", rx, bullet_y, rw, Inches(0.3),
                 font_size=10, font_color=GREEN, bold=True)
        bullet_y += Inches(0.3)
        for s in strengths[:2]:
            add_text(slide, f"✅ {s}", rx, bullet_y, rw, Inches(0.32),
                     font_size=9, font_color=DGRAY)
            bullet_y += Inches(0.3)

    if weaknesses:
        add_text(slide, "Friction / Points faibles", rx, bullet_y, rw, Inches(0.3),
                 font_size=10, font_color=RED_CLR, bold=True)
        bullet_y += Inches(0.3)
        for w in weaknesses[:2]:
            add_text(slide, f"⚠️ {w}", rx, bullet_y, rw, Inches(0.32),
                     font_size=9, font_color=DGRAY)
            bullet_y += Inches(0.3)

    # Verbatim wording
    key_quote = brand_data["evaluation"].get("key_quote")
    if key_quote:
        add_rect(slide, rx, bullet_y + Inches(0.05), rw, Inches(0.55), LGRAY)
        add_rect(slide, rx, bullet_y + Inches(0.05), Inches(0.06), Inches(0.55), ORANGE)
        add_text(slide, f'"{key_quote}"',
                 rx + Inches(0.1), bullet_y + Inches(0.07), rw - Inches(0.15), Inches(0.5),
                 font_size=9, font_color=MARINE, italic=True)

    # Score global en bas
    add_rect(slide, rx, Inches(6.7), rw, Inches(0.65), MARINE)
    add_text(slide, "Score global", rx + Inches(0.1), Inches(6.72), Inches(2), Inches(0.35),
             font_size=10, font_color=WHITE)
    score_txt = f"{'⭐' * int(round(score_g))}  {score_g:.1f} / 5"
    add_text(slide, score_txt, rx + Inches(1.5), Inches(6.72), rw - Inches(1.6), Inches(0.5),
             font_size=13, font_color=ORANGE, bold=True)


# ─── Slide 12 : Tableau synthèse ──────────────────────────────────────

def slide_synthese(prs, brands_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.05), MARINE)
    add_logo(slide, left=Inches(0.2), top=Inches(0.2), height=Inches(0.6))
    add_text(slide, "Tableau de Synthèse Comparatif", Inches(1.6), Inches(0.25), Inches(10), Inches(0.6),
             font_size=24, font_color=WHITE, bold=True)

    # Trier par score global décroissant
    sorted_brands = sorted(brands_data, key=lambda b: b["evaluation"].get("score_global") or 0, reverse=True)

    headers = ["Marque", "Remise", "Fréquences", "Livraison", "Sans engag.", "Score com.", "Score UX", "Score global"]
    col_widths = [Inches(1.8), Inches(1.2), Inches(1.5), Inches(1.3), Inches(1.1), Inches(1.15), Inches(1.0), Inches(1.4)]

    # Header tableau
    x = Inches(0.2)
    for col_i, (header, w) in enumerate(zip(headers, col_widths)):
        add_rect(slide, x, Inches(1.15), w, Inches(0.42), MARINE)
        add_text(slide, header, x + Inches(0.05), Inches(1.18), w - Inches(0.05), Inches(0.38),
                 font_size=9, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += w

    # Lignes
    for row_i, brand in enumerate(sorted_brands):
        cm = brand.get("commercial_model", {})
        ev = brand["evaluation"]
        score_g = ev.get("score_global") or 0
        score_c = ev.get("score_commercial") or 0
        score_u = ev.get("score_ux") or 0
        has_sub = brand["has_subscription"]

        discount = f"{cm.get('discount_pct', 0)}%" if has_sub and cm.get('discount_pct') else "—"
        freqs = str(len(cm.get("frequencies", []))) + " opt." if has_sub and cm.get("frequencies") else "—"
        delivery = "✅ Offerte" if has_sub and cm.get("free_delivery") else ("Seuil" if has_sub else "—")
        engagement = "✅" if has_sub and cm.get("min_engagement") in [None, "sans engagement", "aucun"] else ("❌" if not has_sub else "⚠️")

        row_y = Inches(1.57 + row_i * 0.57)
        bg = LGRAY if row_i % 2 == 0 else WHITE

        values = [brand["brand"], discount, freqs, delivery, engagement,
                  f"{score_c:.1f}/5", f"{score_u:.1f}/5", f"{score_g:.1f}/5"]

        x = Inches(0.2)
        for col_i, (val, w) in enumerate(zip(values, col_widths)):
            add_rect(slide, x, row_y, w, Inches(0.52), bg)

            # Colorer la dernière colonne (score global) selon valeur
            if col_i == len(values) - 1:
                if score_g >= 4.0:
                    cell_bg = RGBColor(0xE8, 0xF8, 0xEF)
                    txt_col = GREEN
                elif score_g >= 3.0:
                    cell_bg = RGBColor(0xFF, 0xF8, 0xE6)
                    txt_col = RGBColor(0xC0, 0x7A, 0x00)
                else:
                    cell_bg = RGBColor(0xFF, 0xF0, 0xF0)
                    txt_col = RED_CLR
                add_rect(slide, x, row_y, w, Inches(0.52), cell_bg)
                add_text(slide, val, x + Inches(0.05), row_y + Inches(0.08), w - Inches(0.05), Inches(0.38),
                         font_size=11, font_color=txt_col, bold=True, align=PP_ALIGN.CENTER)
            else:
                add_text(slide, val, x + Inches(0.05), row_y + Inches(0.08), w - Inches(0.05), Inches(0.38),
                         font_size=10, font_color=DGRAY if col_i > 0 else MARINE,
                         bold=(col_i == 0), align=PP_ALIGN.CENTER if col_i > 0 else PP_ALIGN.LEFT)
            x += w

    # Légende
    legend_y = Inches(6.8)
    add_text(slide, "Légende :  ", Inches(0.2), legend_y, Inches(1.0), Inches(0.4),
             font_size=9, font_color=DGRAY, bold=True)
    add_rect(slide, Inches(1.1), legend_y + Inches(0.1), Inches(0.25), Inches(0.25), RGBColor(0xE8, 0xF8, 0xEF))
    add_text(slide, "≥ 4.0 — Excellent", Inches(1.4), legend_y, Inches(2.5), Inches(0.4),
             font_size=9, font_color=DGRAY)
    add_rect(slide, Inches(3.9), legend_y + Inches(0.1), Inches(0.25), Inches(0.25), RGBColor(0xFF, 0xF8, 0xE6))
    add_text(slide, "3.0–4.0 — Bon", Inches(4.2), legend_y, Inches(2.0), Inches(0.4),
             font_size=9, font_color=DGRAY)
    add_rect(slide, Inches(6.2), legend_y + Inches(0.1), Inches(0.25), Inches(0.25), RGBColor(0xFF, 0xF0, 0xF0))
    add_text(slide, "< 3.0 — Insuffisant / Absent", Inches(6.5), legend_y, Inches(3.0), Inches(0.4),
             font_size=9, font_color=DGRAY)


# ─── Slide 13 : Rentabilité ───────────────────────────────────────────

def slide_rentabilite(prs, renta_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.05), MARINE)
    add_logo(slide, left=Inches(0.2), top=Inches(0.2), height=Inches(0.6))
    add_text(slide, "Analyse Rentabilité Logistique", Inches(1.6), Inches(0.25), Inches(10), Inches(0.6),
             font_size=24, font_color=WHITE, bold=True)

    threshold = renta_data.get("recommended_threshold", 12.35)
    scenarios = renta_data.get("narrative_scenarios", [])

    # 3 scenarios en colonnes
    colors_scenario = [AMBER, GREEN, RGBColor(0x1E, 0x8A, 0x44)]
    icons = ["⚠️", "✅", "🔥"]

    for i, (s, col) in enumerate(zip(scenarios, colors_scenario)):
        x = Inches(0.3 + i * 4.3)
        price = s.get("price", 0)
        net = s.get("net_price", 0)
        margin = s.get("gross_margin", 0)
        logis = s.get("logistics_cost", 5)
        contrib = s.get("contribution_nette", 0)
        viable = s.get("is_profitable", False)
        label = s.get("label", f"Panier {price}€")
        icon = icons[i]

        # Carte
        add_rect(slide, x, Inches(1.2), Inches(4.0), Inches(4.5), LGRAY)
        add_rect(slide, x, Inches(1.2), Inches(4.0), Inches(0.08), col)

        add_text(slide, f"{icon}  {label}", x + Inches(0.15), Inches(1.35), Inches(3.7), Inches(0.45),
                 font_size=11, font_color=MARINE, bold=True)

        rows_r = [
            ("Prix affiché", f"{price:.0f}€ TTC"),
            (f"Après remise -10%", f"{net:.2f}€"),
            (f"Marge brute (45%)", f"{margin:.2f}€"),
            ("Livraison offerte", f"−{logis:.0f}€"),
        ]
        for j, (lbl, val) in enumerate(rows_r):
            ry = Inches(1.9 + j * 0.52)
            bg2 = WHITE if j % 2 == 0 else LGRAY
            add_rect(slide, x, ry, Inches(4.0), Inches(0.48), bg2)
            add_text(slide, lbl, x + Inches(0.1), ry + Inches(0.06), Inches(2.5), Inches(0.38),
                     font_size=9, font_color=DGRAY)
            add_text(slide, val, x + Inches(2.5), ry + Inches(0.06), Inches(1.35), Inches(0.38),
                     font_size=11, font_color=MARINE, bold=True, align=PP_ALIGN.RIGHT)

        # Contribution nette
        add_rect(slide, x, Inches(3.98), Inches(4.0), Inches(0.65), col)
        add_text(slide, "Contribution nette", x + Inches(0.1), Inches(4.03), Inches(2.5), Inches(0.5),
                 font_size=10, font_color=WHITE, bold=True)
        add_text(slide, f"{contrib:.2f}€", x + Inches(2.1), Inches(4.0), Inches(1.7), Inches(0.65),
                 font_size=18, font_color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

    # Encadré seuil
    add_rect(slide, Inches(0.3), Inches(5.85), Inches(12.7), Inches(1.1), MARINE)
    add_rect(slide, Inches(0.3), Inches(5.85), Inches(0.12), Inches(1.1), ORANGE)
    add_text(slide, "💡  Seuil recommandé pour livraison offerte rentable (marge 45%, logistique 5€)",
             Inches(0.55), Inches(5.9), Inches(9), Inches(0.4),
             font_size=11, font_color=LGRAY)
    add_text(slide, f"Panier minimum : 25€",
             Inches(9.5), Inches(5.88), Inches(3.5), Inches(0.55),
             font_size=20, font_color=ORANGE, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, "(Calcul : L ÷ (m × (1−d)) = 5 ÷ (0,45 × 0,90) = 12,35€ → arrondi opérationnel : 25€ min recommandé pour sécurité)",
             Inches(0.55), Inches(6.4), Inches(12), Inches(0.4),
             font_size=8.5, font_color=LGRAY, italic=True)


# ─── Slide 14 : Reco Produits ─────────────────────────────────────────

def slide_reco_produits(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.05), MARINE)
    add_logo(slide, left=Inches(0.2), top=Inches(0.2), height=Inches(0.6))
    add_text(slide, "Recommandation — Produits Éligibles au Lancement", Inches(1.6), Inches(0.25), Inches(11), Inches(0.6),
             font_size=22, font_color=WHITE, bold=True)

    zones = [
        {
            "title": "✅  LANCEMENT V1",
            "subtitle": "Candidats naturels à l'abonnement",
            "color": GREEN,
            "bg": RGBColor(0xE8, 0xF8, 0xEF),
            "items": [
                "Créatine monohydrate (sans goût, récurrence mensuelle naturelle)",
                "Oméga-3 (60-90 caps, cycle de 1-2 mois)",
                "Magnésium bisglycinate (sans goût, prise quotidienne)",
                "Collagène hydrolysé neutre (récurrence élevée)",
                "Multivitamines (prise quotidienne, cycle de 30 jours)",
            ],
            "criteria": "Critères : récurrence naturelle mensuelle + pas de variante goût + SKU simple",
        },
        {
            "title": "⚠️  LANCEMENT V2",
            "subtitle": "À traiter avec précautions",
            "color": AMBER,
            "bg": RGBColor(0xFF, 0xF8, 0xE6),
            "items": [
                "Whey (⚠️ forcer 1 seul goût par abonnement + permettre modification avant envoi)",
                "Preworkout (⚠️ cycles d'utilisation variables, pas mensuel systématiquement)",
                "BCAA / EAA aromatisés (idem risque goût)",
            ],
            "criteria": "Action requise : modification de goût possible avant envoi (comme MyProtein)",
        },
        {
            "title": "❌  HORS PÉRIMÈTRE",
            "subtitle": "À exclure au lancement",
            "color": RED_CLR,
            "bg": RGBColor(0xFF, 0xF0, 0xF0),
            "items": [
                "Produits saisonniers (vitamine D hiver, etc.)",
                "Gammes en révision de formule",
                "Produits à faible récurrence (< 1× / 2 mois)",
            ],
            "criteria": "Raison : risque de rupture de formule + faible valeur LTV abonnement",
        },
    ]

    for i, zone in enumerate(zones):
        x = Inches(0.2 + i * 4.35)
        add_rect(slide, x, Inches(1.15), Inches(4.2), Inches(5.8), zone["bg"])
        add_rect(slide, x, Inches(1.15), Inches(4.2), Inches(0.6), zone["color"])
        add_text(slide, zone["title"], x + Inches(0.1), Inches(1.18), Inches(4.0), Inches(0.4),
                 font_size=11, font_color=WHITE, bold=True)
        add_text(slide, zone["subtitle"], x + Inches(0.1), Inches(1.82), Inches(4.0), Inches(0.3),
                 font_size=9, font_color=zone["color"], bold=True)

        for j, item in enumerate(zone["items"]):
            add_text(slide, f"• {item}", x + Inches(0.1), Inches(2.2 + j * 0.55), Inches(4.0), Inches(0.5),
                     font_size=9.5, font_color=DGRAY)

        add_rect(slide, x, Inches(5.8), Inches(4.2), Inches(0.03), zone["color"])
        add_text(slide, zone["criteria"], x + Inches(0.1), Inches(5.88), Inches(4.0), Inches(0.5),
                 font_size=8.5, font_color=DGRAY, italic=True)


# ─── Slide 15 : Reco Modèle Commercial ───────────────────────────────

def slide_reco_modele(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.05), MARINE)
    add_logo(slide, left=Inches(0.2), top=Inches(0.2), height=Inches(0.6))
    add_text(slide, "Recommandation — Modèle Commercial Impulse", Inches(1.6), Inches(0.25), Inches(11), Inches(0.6),
             font_size=22, font_color=WHITE, bold=True)

    params = [
        ("💸", "Remise abonnement", "−10% sur chaque livraison",
         "Aligné sur le marché FR premium (Nutri&Co, Novoma). Évite la dévaluation de marque\nprovoquée par les remises agressives type MyProtein (−30%). Justifié par l'analyse\nde rentabilité : rentable dès 15€ de panier."),
        ("🔄", "Fréquences proposées", "1 mois / 2 mois / 3 mois",
         "3 options = standard de référence (MyProtein, Novoma). Permet d'adapter\nà chaque type de produit (créatine mensuelle, oméga-3 bimestriel)."),
        ("🚚", "Livraison", "Offerte à partir de 25€ de panier",
         "Seuil issu du modèle de rentabilité (marge 45%, logistique 5€).\nSous 25€ : livraison à 4,50€ relais / 6,50€ domicile — cohérent Nutri&Co."),
        ("🤝", "Engagement", "Sans engagement + bonus fidélité à 6 mois",
         "Sans engagement = levier de conversion maximal. Option : offrir un avantage\n(ex: produit offert ou remise +5%) après 6 livraisons consécutives."),
        ("🛒", "Panier minimum", "25€ pour activer l'abonnement avec livraison offerte",
         "Basé sur l'analyse de rentabilité. Exclut les produits seuls trop petits\net oriente vers les packs ou produits >25€ en V1."),
    ]

    for i, (icon, title, value, rationale) in enumerate(params):
        row_y = Inches(1.2 + i * 1.1)
        add_rect(slide, Inches(0.2), row_y, Inches(12.9), Inches(1.0), LGRAY if i % 2 == 0 else WHITE)
        add_rect(slide, Inches(0.2), row_y, Inches(0.07), Inches(1.0), ORANGE)

        # Icône + titre + valeur
        add_text(slide, icon, Inches(0.35), row_y + Inches(0.1), Inches(0.5), Inches(0.7), font_size=20, font_color=MARINE)
        add_text(slide, title, Inches(0.9), row_y + Inches(0.07), Inches(2.5), Inches(0.38),
                 font_size=10, font_color=DGRAY, bold=False)
        add_text(slide, value, Inches(0.9), row_y + Inches(0.42), Inches(3.3), Inches(0.45),
                 font_size=13, font_color=MARINE, bold=True)

        # Justification
        add_rect(slide, Inches(4.4), row_y + Inches(0.08), Inches(8.6), Inches(0.85), WHITE)
        add_text(slide, rationale, Inches(4.5), row_y + Inches(0.08), Inches(8.5), Inches(0.9),
                 font_size=9, font_color=DGRAY)


# ─── Slide 16 : Reco UX ───────────────────────────────────────────────

def slide_reco_ux(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.05), MARINE)
    add_logo(slide, left=Inches(0.2), top=Inches(0.2), height=Inches(0.6))
    add_text(slide, "Recommandation — UX & Positionnement Impulse", Inches(1.6), Inches(0.25), Inches(11), Inches(0.6),
             font_size=22, font_color=WHITE, bold=True)

    # Wording suggéré
    add_text(slide, "💬  Wording suggéré (module fiche produit)", Inches(0.3), Inches(1.2), Inches(8), Inches(0.4),
             font_size=13, font_color=MARINE, bold=True)
    add_rect(slide, Inches(0.3), Inches(1.65), Inches(12.7), Inches(0.9), LGRAY)
    add_rect(slide, Inches(0.3), Inches(1.65), Inches(0.12), Inches(0.9), ORANGE)
    add_text(slide, '"Abonnez-vous & économisez 10% — Livraison offerte dès 25€ — Pause quand vous voulez"',
             Inches(0.55), Inches(1.75), Inches(12.3), Inches(0.65),
             font_size=14, font_color=MARINE, italic=True, bold=False)

    # 3 points d'entrée
    add_text(slide, "📍  Points d'entrée recommandés", Inches(0.3), Inches(2.75), Inches(8), Inches(0.4),
             font_size=13, font_color=MARINE, bold=True)

    entries = [
        ("Fiche produit", "Module au-dessus du bouton 'Ajouter au panier'\nOptions : Achat unique | S'abonner et économiser\nPlacement : AVANT le CTA (best practice MyProtein)"),
        ("Homepage", "Bandeau ou encadré hero : 'Commandez en abonnement\net économisez 10% à chaque livraison'\nTest A/B recommandé vs version sans mise en avant"),
        ("Navigation & email", "• Entrée 'Abonnement' dans le menu principal\n• Email post-achat J+14 : 'Passez à l'abonnement'\n• Email de rappel à J-5 avant épuisement estimé"),
    ]

    for i, (title, body) in enumerate(entries):
        x = Inches(0.3 + i * 4.35)
        add_rect(slide, x, Inches(3.25), Inches(4.1), Inches(3.35), LGRAY)
        add_rect(slide, x, Inches(3.25), Inches(4.1), Inches(0.5), MARINE)
        add_text(slide, title, x + Inches(0.15), Inches(3.3), Inches(3.8), Inches(0.42),
                 font_size=12, font_color=WHITE, bold=True)
        add_text(slide, body, x + Inches(0.15), Inches(3.82), Inches(3.8), Inches(2.5),
                 font_size=10, font_color=DGRAY)

    # Note benchmark
    add_rect(slide, Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.55), LGRAY)
    add_text(slide,
             "📌  Référence benchmark : MyProtein (score UX 4.5/5) place le module AVANT le CTA avec 2 options visuellement distinctes — à reproduire.",
             Inches(0.45), Inches(6.82), Inches(12.3), Inches(0.42),
             font_size=9.5, font_color=MARINE, italic=False)


# ─── Slide 17 : Prochaines étapes ────────────────────────────────────

def slide_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.05), MARINE)
    add_logo(slide, left=Inches(0.2), top=Inches(0.2), height=Inches(0.6))
    add_text(slide, "Prochaines Étapes — Roadmap Lancement", Inches(1.6), Inches(0.25), Inches(11), Inches(0.6),
             font_size=24, font_color=WHITE, bold=True)

    steps = [
        ("Mai 2026", "Choix technologique", [
            "Shopify Subscriptions natif (gratuit, simple)",
            "Recharge / Bold / Loop (plus flexible, payant)",
            "Décision : budget vs. besoin de personnalisation",
        ], ORANGE),
        ("Juin 2026", "Setup technique & sélection SKUs", [
            "Intégration de l'app choisie sur Shopify",
            "Sélection des 5 produits V1 éligibles",
            "Configuration remise 10% + seuil 25€",
        ], MARINE),
        ("Juillet 2026", "Lancement bêta fermé", [
            "Invitation clients existants (newsletter)",
            "Test UX module fiche produit",
            "Mesure taux de souscription & churn J30",
        ], RGBColor(0x2D, 0xA8, 0x6A)),
        ("Août 2026", "Lancement public + push marketing", [
            "Ouverture à tous les visiteurs",
            "Email campaign + social media",
            "Suivi KPIs : LTV abonnés vs one-shot",
        ], RGBColor(0x5A, 0x6F, 0xC0)),
    ]

    for i, (date, title, bullets, color) in enumerate(steps):
        x = Inches(0.2 + i * 3.28)
        add_rect(slide, x, Inches(1.2), Inches(3.1), Inches(5.5), LGRAY)
        add_rect(slide, x, Inches(1.2), Inches(3.1), Inches(0.55), color)
        add_text(slide, date, x + Inches(0.1), Inches(1.23), Inches(2.9), Inches(0.28),
                 font_size=9, font_color=WHITE)
        add_text(slide, title, x + Inches(0.1), Inches(1.52), Inches(2.9), Inches(0.42),
                 font_size=12, font_color=WHITE, bold=True)
        for j, b in enumerate(bullets):
            add_text(slide, f"• {b}", x + Inches(0.1), Inches(1.9 + j * 0.58), Inches(2.9), Inches(0.5),
                     font_size=9.5, font_color=DGRAY)

    # KPIs à suivre
    add_rect(slide, Inches(0.2), Inches(6.85), Inches(12.9), Inches(0.5), MARINE)
    add_text(slide,
             "KPIs post-lancement :  Taux de souscription  •  Churn mensuel (cible <10%)  •  LTV abonnés vs one-shot  •  Panier moyen abonnés  •  Taux de pause vs résiliation",
             Inches(0.4), Inches(6.91), Inches(12.5), Inches(0.4),
             font_size=9.5, font_color=WHITE, align=PP_ALIGN.CENTER)


# ─── Main ─────────────────────────────────────────────────────────────

def main():
    print("🎨 Génération du deck PowerPoint...\n")

    # Chargement des données
    with open(DATA_PATH, encoding="utf-8") as f:
        data = json.load(f)
    brands_data = data["brands"]

    with open(RENTA_PATH, encoding="utf-8") as f:
        renta_data = json.load(f)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("  Slide 1 — Cover")
    slide_cover(prs)

    print("  Slide 2 — Contexte & Objectif")
    slide_contexte(prs)

    print("  Slide 3 — Méthodologie")
    slide_methodologie(prs)

    # Ordre des marques pour les slides (du plus pertinent pour Impulse au plus éloigné)
    brand_order = ["Nutri&Co", "Nutrimuscle", "Novoma", "Cuure", "MyProtein", "Decathlon", "Aqeelab", "Nutripure"]

    for brand_name in brand_order:
        brand = next((b for b in brands_data if b["brand"] == brand_name), None)
        if brand:
            print(f"  Slide — {brand_name}")
            slide_brand(prs, brand)

    print("  Slide 12 — Tableau synthèse comparatif")
    slide_synthese(prs, brands_data)

    print("  Slide 13 — Analyse rentabilité logistique")
    slide_rentabilite(prs, renta_data)

    print("  Slide 14 — Recommandation Produits éligibles")
    slide_reco_produits(prs)

    print("  Slide 15 — Recommandation Modèle commercial")
    slide_reco_modele(prs)

    print("  Slide 16 — Recommandation UX & Positionnement")
    slide_reco_ux(prs)

    print("  Slide 17 — Prochaines étapes")
    slide_next_steps(prs)

    prs.save(str(OUTPUT_PATH))
    size = OUTPUT_PATH.stat().st_size
    print(f"\n✅ Deck généré → {OUTPUT_PATH}")
    print(f"   Taille : {size:,} bytes ({size/1024/1024:.1f} MB)")
    print(f"   Slides : {len(prs.slides)}")


if __name__ == "__main__":
    main()
