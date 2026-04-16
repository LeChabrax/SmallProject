"""
build_deck.py — Génère le deck PPTX final de la mission benchmark abonnement
pour Impulse Nutrition, au format Havea (photo hero + palette mint + layouts denses).

Source : 8 benchmark/marques/<slug>/data.json + insights narratifs depuis
synthesis_table.md, recommendation_impulse.md, profitability_section.md.

Structure 16 slides :
  01 — Cover
  02 — Méthodologie & périmètre
  03–10 — Benchmark par marque (8 slides, 1 par marque)
  11 — Tableau de synthèse comparatif
  12 — Analyse rentabilité / logistique
  13 — Recommandation · produits éligibles
  14 — Recommandation · modèle commercial
  15 — Recommandation · UX et positionnement
  16 — Prochaines étapes
  (17 — Closing optionnel)

Output : benchmark/deck_impulse_abonnement_v1.pptx
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches

# Imports locaux
import sys
sys.path.insert(0, str(Path(__file__).resolve().parent))
import deck_style as ds  # noqa: E402

ROOT = Path(__file__).resolve().parents[1]
MARQUES = ROOT / "marques"
ASSETS = ROOT / "deck" / "assets"
OUT = ROOT / "deck_impulse_abonnement_v1.pptx"

# Ordre d'affichage des marques dans le deck (= classement global, Nutripure à la fin comme contrôle)
BRAND_ORDER = [
    "aqeelab",
    "nutriandco",
    "nutrimuscle",
    "novoma",
    "cuure",
    "myprotein",
    "decathlon",
    "nutripure",
]

# Screenshot principal à utiliser pour chaque marque sur la slide dédiée
BRAND_SCREENSHOT = {
    "nutriandco": "abonnement_page.png",
    "nutrimuscle": "abonnement_page.png",
    "novoma": "abonnement_page.png",
    "nutripure": "homepage.png",
    "cuure": "quiz.png",
    "decathlon": "homepage.png",
    "aqeelab": "product_sub.png",
    "myprotein": "product_sub.png",
}

NAV_SECTIONS = [
    "Intro",
    "Méthodo",
    "Benchmark",
    "Synthèse",
    "Rentabilité",
    "Recommandation",
    "Next",
]


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank_layout)


# ─── Slide 1 — Cover ───────────────────────────────────────────────────

def add_slide_cover(prs):
    slide = blank_slide(prs)
    ds.add_cover_hero(
        slide,
        photo_path=ASSETS / "hero_1.jpeg",
        title="BENCHMARK\nABONNEMENT",
        subtitle="Modèles du marché FR · recommandation Impulse",
        week_label="Comité projet · 2026-04-15",
    )
    return slide


# ─── Slide 2 — Méthodologie & périmètre ────────────────────────────────

def add_slide_methodology(prs, page_num, total):
    slide = blank_slide(prs)
    ds.add_nav_breadcrumb(slide, NAV_SECTIONS, active_idx=1)
    ds.add_editorial_title(
        slide,
        "Méthodologie · 8 marques auditées en live, 4 axes d'analyse, sources tracées",
        size=18,
    )

    # Panneau gauche : Périmètre
    ds.add_textbox(
        slide,
        "PÉRIMÈTRE",
        Inches(0.4),
        Inches(1.5),
        Inches(4),
        Inches(0.25),
        font=ds.BODY_FONT,
        size=9,
        bold=True,
        color=ds.MINT,
    )
    brands_list = [
        "Nutri&Co · référence premium FR la plus proche d'Impulse",
        "Nutrimuscle · leader volume muscu FR",
        "Novoma · santé-sport, bon modèle abo présumé",
        "Nutripure · contrôle sans abonnement",
        "Cuure · modèle 100% abo personnalisé",
        "Decathlon · grande distribution omnicanale",
        "Aqeelab · pure player FR récent premium",
        "MyProtein · 'subscribe & save' international mature",
    ]
    ds.add_bullets(slide, brands_list, 0.4, 1.8, 4.5, 3, size=9)

    # Panneau droit : Structure
    ds.add_textbox(
        slide,
        "STRUCTURE D'ANALYSE (BRIEF)",
        Inches(5.2),
        Inches(1.5),
        Inches(4.5),
        Inches(0.25),
        font=ds.BODY_FONT,
        size=9,
        bold=True,
        color=ds.MINT,
    )
    axes = [
        "1. Modèle commercial — remise, fréquences, livraison, engagement, stacking",
        "2. Périmètre produit — scope, exclusions, ticket moyen",
        "3. UX / page produit — wording, placement, page dédiée, mention home/nav",
        "4. Évaluation — forces, faiblesses, scores /5 (UX · Offre · Pertinence · Global)",
    ]
    ds.add_bullets(slide, axes, 5.2, 1.8, 4.5, 2.5, size=9)

    # Sources
    ds.add_textbox(
        slide,
        "SOURCES",
        Inches(5.2),
        Inches(3.8),
        Inches(4.5),
        Inches(0.25),
        font=ds.BODY_FONT,
        size=9,
        bold=True,
        color=ds.MINT,
    )
    ds.add_textbox(
        slide,
        "Pour chaque marque : pages FAQ officielles + Zendesk Help Center + vérifs WebFetch 2026-04-15 + "
        "Trustpilot + blogs officiels. Chaque citation est tracée dans benchmark/marques/<brand>/sources.md. "
        "Statut audit : 8/8 verified_live.",
        Inches(5.2),
        Inches(4.1),
        Inches(4.5),
        Inches(1.0),
        font=ds.BODY_FONT,
        size=9,
        color=ds.GREY_700,
        italic=True,
        line_spacing=1.3,
    )

    ds.add_havea_footer(slide, page_num, total)


# ─── Slides 3-10 — Benchmark par marque ────────────────────────────────

def _short(text, max_len):
    if not text:
        return ""
    text = str(text)
    return text if len(text) <= max_len else text[: max_len - 1] + "…"


def add_slide_brand(prs, brand_data, page_num, total):
    slide = blank_slide(prs)
    ds.add_nav_breadcrumb(slide, NAV_SECTIONS, active_idx=2)

    brand_name = brand_data.get("brand", "—")
    has_sub = brand_data.get("has_subscription", False)

    # Titre éditorial : nom + verdict en 1 ligne
    mc = brand_data.get("modele_commercial", {}) or {}
    remise = mc.get("remise_pct")
    if remise is None and not has_sub:
        remise_label = "Pas d'abonnement"
    elif isinstance(remise, (int, float)):
        remise_label = f"−{remise}%"
    else:
        remise_label = f"−{remise}" if remise else "—"

    ev = brand_data.get("evaluation", {}) or {}
    global_score = ev.get("score_global")
    score_label = f"· Note globale {global_score}/5" if global_score else ""

    editorial = f"{brand_name} · {remise_label}  {score_label}"
    ds.add_editorial_title(slide, editorial, size=18)

    # ─── Colonne gauche : screenshot ───
    slug = brand_data.get("slug")
    screenshot_file = BRAND_SCREENSHOT.get(slug, "homepage.png")
    screenshot_path = MARQUES / slug / "screenshots" / screenshot_file
    caption = f"{brand_data.get('url', '')} · {screenshot_file}"
    ds.add_screenshot_with_caption(
        slide,
        image_path=screenshot_path,
        x=0.4,
        y=1.55,
        w=4.4,
        h=3.5,
        caption=caption,
    )

    # ─── Colonne droite : 4 rubriques + scorecard ───
    RX = 5.1
    RW = 4.5

    # Rubrique 1 — Modèle commercial
    ds.add_textbox(
        slide,
        "01 · MODÈLE COMMERCIAL",
        Inches(RX),
        Inches(1.55),
        Inches(RW),
        Inches(0.2),
        font=ds.BODY_FONT,
        size=8,
        bold=True,
        color=ds.MINT,
    )

    freqs = mc.get("frequencies", [])
    livraison = mc.get("livraison_abo", {}) or {}
    if livraison.get("offerte") is True and livraison.get("seuil_eur") == 0:
        liv_txt = "livraison offerte sans seuil"
    elif livraison.get("offerte") == "relais uniquement":
        liv_txt = "livraison offerte en relais uniquement"
    elif livraison.get("offerte") is True and livraison.get("seuil_eur"):
        liv_txt = f"livraison offerte dès {livraison.get('seuil_eur')}€"
    elif livraison.get("offerte") is False:
        liv_txt = f"livraison stepped (dès {livraison.get('seuil_eur', '?')}€)"
    else:
        liv_txt = "livraison : non documenté"

    eng = mc.get("engagement_min", "")
    if eng and "Aucun" in eng:
        eng_txt = "sans engagement"
    elif eng and "2 mois" in eng:
        eng_txt = "engagement 2 mois min"
    elif eng and "partiel" in eng.lower():
        eng_txt = "engagement partiel (friction)"
    elif not has_sub:
        eng_txt = "N/A (pas d'abo)"
    else:
        eng_txt = _short(eng, 40)

    line_1 = f"Remise : {remise_label}  ·  Fréquences : {len(freqs)} ({', '.join(freqs)[:30]})" if freqs else f"Remise : {remise_label}"
    line_2 = f"{liv_txt.capitalize()}  ·  {eng_txt.capitalize()}"
    ds.add_textbox(
        slide,
        line_1,
        Inches(RX),
        Inches(1.78),
        Inches(RW),
        Inches(0.2),
        font=ds.BODY_FONT,
        size=9,
        color=ds.INK,
    )
    ds.add_textbox(
        slide,
        line_2,
        Inches(RX),
        Inches(1.98),
        Inches(RW),
        Inches(0.2),
        font=ds.BODY_FONT,
        size=9,
        color=ds.GREY_700,
    )

    # Rubrique 2 — Périmètre produit
    ds.add_textbox(
        slide,
        "02 · PÉRIMÈTRE PRODUIT",
        Inches(RX),
        Inches(2.3),
        Inches(RW),
        Inches(0.2),
        font=ds.BODY_FONT,
        size=8,
        bold=True,
        color=ds.MINT,
    )
    pp = brand_data.get("perimetre_produit", {}) or {}
    scope = pp.get("scope") or "—"
    ticket = pp.get("ticket_moyen_estime_eur")
    ticket_txt = f" · Panier moyen ~{ticket}€" if ticket else ""
    scope_line = _short(scope, 90) + ticket_txt
    ds.add_textbox(
        slide,
        scope_line,
        Inches(RX),
        Inches(2.53),
        Inches(RW),
        Inches(0.4),
        font=ds.BODY_FONT,
        size=9,
        color=ds.GREY_700,
        line_spacing=1.25,
    )

    # Rubrique 3 — UX
    ds.add_textbox(
        slide,
        "03 · UX / PAGE PRODUIT",
        Inches(RX),
        Inches(3.0),
        Inches(RW),
        Inches(0.2),
        font=ds.BODY_FONT,
        size=8,
        bold=True,
        color=ds.MINT,
    )
    ux = brand_data.get("ux_page_produit", {}) or {}
    cta = _short(ux.get("wording_cta"), 30)
    dedicated = "✓ page dédiée" if ux.get("dedicated_page") else "✗ pas de page dédiée"
    home = "✓ home" if ux.get("homepage_mention") and ux.get("homepage_mention") is not False else "✗ home"
    nav = "✓ nav" if ux.get("navigation_mention") and ux.get("navigation_mention") is not False else "✗ nav"
    ux_line = f'CTA : "{cta}"  ·  {dedicated}  ·  {home}  ·  {nav}'
    ds.add_textbox(
        slide,
        ux_line,
        Inches(RX),
        Inches(3.23),
        Inches(RW),
        Inches(0.25),
        font=ds.BODY_FONT,
        size=9,
        color=ds.GREY_700,
    )

    # Rubrique 4 — Évaluation : 1 force + 1 faiblesse + highlight
    ds.add_textbox(
        slide,
        "04 · ÉVALUATION",
        Inches(RX),
        Inches(3.6),
        Inches(RW),
        Inches(0.2),
        font=ds.BODY_FONT,
        size=8,
        bold=True,
        color=ds.MINT,
    )
    strengths = ev.get("strengths", [])
    weaknesses = ev.get("weaknesses", [])
    force = _short(strengths[0], 85) if strengths else "—"
    faiblesse = _short(weaknesses[0], 85) if weaknesses else "—"
    ds.add_multiline(
        slide,
        [
            ("+  " + force, 9, False, ds.MINT, False),
            ("−  " + faiblesse, 9, False, ds.SALMON, False),
        ],
        Inches(RX),
        Inches(3.83),
        Inches(RW),
        Inches(0.8),
        line_spacing=1.3,
    )

    # Scorecard
    ds.add_scorecard(
        slide,
        x=5.1,
        y=4.55,
        scores={
            "ux": ev.get("score_ux"),
            "offre": ev.get("score_offre_commerciale"),
            "pertinence": ev.get("score_pertinence_vs_impulse"),
            "global": ev.get("score_global"),
        },
        cell_w=1.05,
        cell_h=0.75,
        gap=0.05,
    )

    ds.add_havea_footer(slide, page_num, total)


# ─── Slide 11 — Tableau de synthèse comparatif ─────────────────────────

def fmt_remise_cell(d):
    mc = d.get("modele_commercial", {}) or {}
    pct = mc.get("remise_pct")
    if pct is None:
        return "—"
    if isinstance(pct, (int, float)):
        return f"−{pct}%"
    return str(pct)[:14]


def fmt_liv_cell(d):
    liv = d.get("modele_commercial", {}).get("livraison_abo", {}) or {}
    o = liv.get("offerte")
    s = liv.get("seuil_eur")
    if o is True and s == 0:
        return "✓ sans seuil"
    if o == "relais uniquement":
        return "✓ relais only"
    if o is True and s:
        return f"✓ dès {s}€"
    if o is False:
        return f"stepped {s}€" if s else "stepped"
    return "—"


def fmt_eng_cell(d):
    eng = d.get("modele_commercial", {}).get("engagement_min", "") or ""
    if "Aucun" in eng or eng == "":
        return "Aucun"
    if "2 mois" in eng:
        return "2 mois min"
    if "partiel" in eng.lower() or "friction" in eng.lower():
        return "Friction"
    return eng[:12]


def add_slide_synthesis(prs, brands, page_num, total):
    slide = blank_slide(prs)
    ds.add_nav_breadcrumb(slide, NAV_SECTIONS, active_idx=3)
    ds.add_editorial_title(
        slide,
        "Tableau de synthèse · 8 marques × 7 dimensions + scores",
        size=18,
    )

    headers = ["Marque", "Remise", "Fréq.", "Livraison", "Engag.", "UX", "Offre", "Pert.", "Global"]
    col_widths = [1.6, 0.9, 0.55, 1.35, 0.9, 0.55, 0.65, 0.75, 0.75]

    rows = []
    for slug in BRAND_ORDER:
        b = brands[slug]
        ev = b.get("evaluation", {}) or {}
        mc = b.get("modele_commercial", {}) or {}
        freq_count = len(mc.get("frequencies", []))
        rows.append([
            b.get("brand", "—")[:22],
            fmt_remise_cell(b),
            str(freq_count) if freq_count else "—",
            fmt_liv_cell(b),
            fmt_eng_cell(b),
            _score(ev.get("score_ux")),
            _score(ev.get("score_offre_commerciale")),
            _score(ev.get("score_pertinence_vs_impulse")),
            _score(ev.get("score_global")),
        ])

    ds.add_data_table(
        slide,
        headers,
        rows,
        x=0.4,
        y=1.5,
        w=9.2,
        col_widths=col_widths,
        row_h=0.28,
        header_h=0.28,
        font_size=9,
    )

    # Légende + insight
    ds.add_insight_box(
        slide,
        x=0.4,
        y=4.85,
        w=9.2,
        h=0.65,
        title="Lecture",
        body="Top 3 par note globale : Aqeelab (4,3) · Nutri&Co (4,0) · Nutrimuscle (3,8). "
        "Deux patterns antagonistes : acquisition agressive (livraison offerte + remise élevée) vs rentabilité prudente (livraison stepped).",
    )

    ds.add_havea_footer(slide, page_num, total)


def _score(v):
    if v is None:
        return "—"
    return f"{v}"


# ─── Slide 12 — Rentabilité ────────────────────────────────────────────

def add_slide_profitability(prs, page_num, total):
    slide = blank_slide(prs)
    ds.add_nav_breadcrumb(slide, NAV_SECTIONS, active_idx=4)
    ds.add_editorial_title(
        slide,
        "Rentabilité · 3 scénarios du brief modélisés avec COGS réels Impulse, remise −15%",
        size=17,
    )

    # 3 P&L côte à côte
    pnls = [
        {
            "title": "Mono-produit ~15€",
            "sub": "Magnésium Bisglycinate 15,90€",
            "rows": [
                ("Prix retail HT", "13,25 €"),
                ("Remise −15%", "−1,99 €"),
                ("PVM net HT", "11,26 €"),
                ("COGS HT", "−3,16 €"),
                ("Coût log HT", "−4,92 €"),
            ],
            "total": "3,18 €",
            "label": "Contribution nette",
            "pct": "24%",
            "verdict": "Rentable · tendu",
        },
        {
            "title": "Panier ~35€",
            "sub": "Magnésium + Multivit 34,80€",
            "rows": [
                ("Panier HT", "29,00 €"),
                ("Remise −15%", "−4,35 €"),
                ("PVM net HT", "24,65 €"),
                ("COGS HT", "−4,38 €"),
                ("Coût log HT", "−4,92 €"),
            ],
            "total": "15,35 €",
            "label": "Contribution nette",
            "pct": "53%",
            "verdict": "Très rentable",
        },
        {
            "title": "Pack 3 SKUs ~50€",
            "sub": "Magnésium + Multivit + Créatine 47,70€",
            "rows": [
                ("Panier HT", "39,75 €"),
                ("Remise −15%", "−5,96 €"),
                ("PVM net HT", "33,79 €"),
                ("COGS HT", "−8,04 €"),
                ("Coût log HT", "−4,92 €"),
            ],
            "total": "20,83 €",
            "label": "Contribution nette",
            "pct": "52%",
            "verdict": "Idéal · cible upsell",
        },
    ]

    card_w = 3.0
    gap = 0.15
    start_x = 0.4
    for i, p in enumerate(pnls):
        x = start_x + i * (card_w + gap)
        # Fond carte
        ds.add_rect(slide, Inches(x), Inches(1.5), Inches(card_w), Inches(2.85), fill_color=ds.GREY_100)
        # Accent mint top
        ds.add_rect(slide, Inches(x), Inches(1.5), Inches(card_w), Inches(0.08), fill_color=ds.MINT)
        # Titre
        ds.add_textbox(
            slide,
            p["title"].upper(),
            Inches(x + 0.2),
            Inches(1.65),
            Inches(card_w - 0.4),
            Inches(0.25),
            font=ds.BODY_FONT,
            size=9,
            bold=True,
            color=ds.MINT,
        )
        ds.add_textbox(
            slide,
            p["sub"],
            Inches(x + 0.2),
            Inches(1.85),
            Inches(card_w - 0.4),
            Inches(0.2),
            font=ds.BODY_FONT,
            size=8,
            italic=True,
            color=ds.GREY_500,
        )
        # Rows
        y = 2.12
        for label, val in p["rows"]:
            ds.add_textbox(
                slide,
                label,
                Inches(x + 0.2),
                Inches(y),
                Inches(1.6),
                Inches(0.2),
                font=ds.BODY_FONT,
                size=9,
                color=ds.INK,
            )
            ds.add_textbox(
                slide,
                val,
                Inches(x + 1.4),
                Inches(y),
                Inches(card_w - 1.6),
                Inches(0.2),
                font=ds.MONO_FONT,
                size=9,
                color=ds.INK,
                align=PP_ALIGN.RIGHT,
            )
            y += 0.24
        # Separator
        ds.add_rect(
            slide,
            Inches(x + 0.2),
            Inches(y + 0.04),
            Inches(card_w - 0.4),
            Inches(0.015),
            fill_color=ds.INK,
        )
        # Total
        ds.add_textbox(
            slide,
            p["label"],
            Inches(x + 0.2),
            Inches(y + 0.12),
            Inches(1.6),
            Inches(0.22),
            font=ds.BODY_FONT,
            size=9,
            bold=True,
            color=ds.INK,
        )
        ds.add_textbox(
            slide,
            p["total"],
            Inches(x + 1.4),
            Inches(y + 0.12),
            Inches(card_w - 1.6),
            Inches(0.22),
            font=ds.MONO_FONT,
            size=12,
            bold=True,
            color=ds.MINT,
            align=PP_ALIGN.RIGHT,
        )
        # Verdict
        ds.add_textbox(
            slide,
            f"{p['verdict']}  ·  {p['pct']}",
            Inches(x + 0.2),
            Inches(y + 0.4),
            Inches(card_w - 0.4),
            Inches(0.2),
            font=ds.BODY_FONT,
            size=8,
            italic=True,
            color=ds.GREY_700,
        )

    # Insight bottom
    ds.add_insight_box(
        slide,
        x=0.4,
        y=4.5,
        w=9.2,
        h=0.95,
        title="Conclusion rentabilité",
        body="Marges brutes Impulse exceptionnelles (62-92% HT). Coût log compétitif (3,25€ relais / 4,92€ domicile HT, bas de fourchette marché). "
        "Break-even à 39 abonnés (app) · LTV baseline 378€ · ratio LTV/CAC très sain.",
    )

    ds.add_havea_footer(slide, page_num, total)


# ─── Slides 13, 14, 15 — Recommandation Impulse ─────────────────────────

def add_slide_reco_products(prs, page_num, total):
    slide = blank_slide(prs)
    ds.add_nav_breadcrumb(slide, NAV_SECTIONS, active_idx=5)
    ds.add_editorial_title(
        slide,
        "Recommandation 1/3 · 6 SKUs MVP à consommation naturellement récurrente",
        size=17,
    )

    # Colonne Go (6 SKUs)
    ds.add_textbox(
        slide,
        "✓  MVP PHASE 1 · 6 SKUs",
        Inches(0.4),
        Inches(1.5),
        Inches(4.5),
        Inches(0.25),
        font=ds.BODY_FONT,
        size=9,
        bold=True,
        color=ds.MINT,
    )
    go_skus = [
        "Magnésium Bisglycinate · santé minéraux · 1-2 mois/pot",
        "Oméga-3 EPAX® · santé acides gras · quotidien · SKU phare",
        "Multivitamines · santé vitamines · panier d'entrée naturel",
        "Pré-Probiotiques · santé digestion · cure 2-3 mois",
        "Sommeil + · santé adaptogènes · quotidien le soir",
        "Créatine Monohydrate · sport perf · 5g/jour · acquisition sportive",
    ]
    ds.add_bullets(slide, go_skus, 0.4, 1.8, 4.5, 3.5, size=10)

    # Colonne Maybe + Exclusions
    ds.add_textbox(
        slide,
        "⚠  PHASE 2 · avec précaution",
        Inches(5.2),
        Inches(1.5),
        Inches(4.5),
        Inches(0.25),
        font=ds.BODY_FONT,
        size=9,
        bold=True,
        color=ds.BEIGE_GOLD,
    )
    maybe = [
        "Whey · lock goût à la souscription (bouton changer dans espace abonné)",
        "Collagène Marin · selon format et goût après validation Phase 1",
    ]
    ds.add_bullets(slide, maybe, 5.2, 1.8, 4.5, 1.0, size=10, marker_color=ds.BEIGE_GOLD)

    ds.add_textbox(
        slide,
        "✗  EXCLUS DU LANCEMENT",
        Inches(5.2),
        Inches(3.1),
        Inches(4.5),
        Inches(0.25),
        font=ds.BODY_FONT,
        size=9,
        bold=True,
        color=ds.SALMON,
    )
    excluded = [
        "Packs / coffrets / bundles · composition figée",
        "Accessoires · pas de récurrence naturelle",
        "Snacks · consommation irrégulière",
        "Électrolytes effervescents · événementiel",
        "Produits en rupture de formule ou saisonniers",
    ]
    ds.add_bullets(slide, excluded, 5.2, 3.4, 4.5, 2, size=10, marker_color=ds.SALMON)

    ds.add_havea_footer(slide, page_num, total)


def add_slide_reco_commercial(prs, page_num, total):
    slide = blank_slide(prs)
    ds.add_nav_breadcrumb(slide, NAV_SECTIONS, active_idx=5)
    ds.add_editorial_title(
        slide,
        "Recommandation 2/3 · Modèle commercial −15% flat, sans engagement, livraison offerte dès 35€",
        size=16,
    )

    # 5 cartes verticales
    cards = [
        ("REMISE", "−15%", "flat · aligné Nutri&Co et Nutrimuscle", "Sous Aqeelab (−20%) mais compensé par zéro engagement."),
        ("FRÉQUENCES", "1 / 2 / 3", "mois · point de différenciation", "Le '2 mois' manque chez Nutri&Co — sweet spot pot 60 gélules."),
        ("LIVRAISON", "dès 35€", "domicile offerte · seuil incitatif", "Modèle hybride protège marge mono-produit, pousse à l'upsell."),
        ("ENGAGEMENT", "Aucun", "résiliation 1 clic · Loi Chatel", "Leçon Aqeelab (2 mois = friction) + Cuure (2 temps = anti-pattern)."),
        ("STACKING", "Partiel", "welcome OK · promos non cumul", "Approche Nutrimuscle : booste acquisition sans diluer la marge abo."),
    ]
    card_w = 1.85
    gap = 0.1
    start_x = 0.4
    for i, (label, value, sub, note) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        ds.add_rect(slide, Inches(x), Inches(1.5), Inches(card_w), Inches(3.0), fill_color=ds.GREY_100)
        ds.add_rect(slide, Inches(x), Inches(1.5), Inches(card_w), Inches(0.06), fill_color=ds.MINT)
        ds.add_textbox(
            slide,
            label,
            Inches(x + 0.15),
            Inches(1.65),
            Inches(card_w - 0.3),
            Inches(0.2),
            font=ds.BODY_FONT,
            size=8,
            bold=True,
            color=ds.MINT,
        )
        ds.add_textbox(
            slide,
            value,
            Inches(x + 0.15),
            Inches(1.9),
            Inches(card_w - 0.3),
            Inches(0.6),
            font=ds.DISPLAY_FONT,
            size=22,
            color=ds.INK,
            bold=False,
            align=PP_ALIGN.CENTER,
        )
        ds.add_textbox(
            slide,
            sub,
            Inches(x + 0.15),
            Inches(2.6),
            Inches(card_w - 0.3),
            Inches(0.35),
            font=ds.BODY_FONT,
            size=8,
            italic=True,
            color=ds.GREY_700,
            align=PP_ALIGN.CENTER,
            line_spacing=1.2,
        )
        ds.add_textbox(
            slide,
            note,
            Inches(x + 0.15),
            Inches(3.1),
            Inches(card_w - 0.3),
            Inches(1.2),
            font=ds.BODY_FONT,
            size=8,
            color=ds.INK,
            line_spacing=1.3,
        )

    # Insight bottom
    ds.add_insight_box(
        slide,
        x=0.4,
        y=4.75,
        w=9.2,
        h=0.75,
        title="Verdict commercial",
        body="−15% flat est prudent (Impulse peut aller jusqu'à −30% en restant rentable) et aligné sur la référence marché. "
        "Le seuil 35€ est une incitation à l'upsell plutôt qu'une protection marge.",
    )

    ds.add_havea_footer(slide, page_num, total)


def add_slide_reco_ux(prs, page_num, total):
    slide = blank_slide(prs)
    ds.add_nav_breadcrumb(slide, NAV_SECTIONS, active_idx=5)
    ds.add_editorial_title(
        slide,
        "Recommandation 3/3 · UX persuasive, visibilité maximale, anti-churn éprouvé",
        size=17,
    )

    # Colonne À copier (4 patterns)
    ds.add_textbox(
        slide,
        "À COPIER · PATTERNS GAGNANTS DU PANEL",
        Inches(0.4),
        Inches(1.5),
        Inches(5),
        Inches(0.25),
        font=ds.BODY_FONT,
        size=9,
        bold=True,
        color=ds.MINT,
    )
    to_copy = [
        'Wording CTA  "S\'abonner & économiser"  · source Aqeelab',
        'Affichage du montant économisé en euros sur la fiche  · source Aqeelab (pattern unique)',
        'Page dédiée /pages/abonnement avec pitch éditorial santé  · source Novoma',
        'Mention homepage + navigation + newsletter  · source Nutrimuscle',
        'Email pré-expédition 48h avec fenêtre de modification  · source Nutri&Co (anti-churn)',
    ]
    ds.add_bullets(slide, to_copy, 0.4, 1.8, 9.2, 2.5, size=10, marker_color=ds.MINT)

    # Colonne À éviter
    ds.add_textbox(
        slide,
        "À ÉVITER · ANTI-PATTERNS DOCUMENTÉS",
        Inches(0.4),
        Inches(3.5),
        Inches(5),
        Inches(0.25),
        font=ds.BODY_FONT,
        size=9,
        bold=True,
        color=ds.SALMON,
    )
    to_avoid = [
        "Sous-domaine séparé  · Decathlon (friction redirect)",
        '"1 produit = 1 abo"  · Decathlon (mauvaise UX multi-abos)',
        "Résiliation en 2 temps (suspension vs email)  · Cuure (dark pattern)",
        "Engagement minimum forcé  · Aqeelab (friction Loi Chatel)",
        "Prix de base instables  · MyProtein (remise illisible)",
    ]
    ds.add_bullets(slide, to_avoid, 0.4, 3.8, 9.2, 2.5, size=10, marker_color=ds.SALMON, marker="✗")

    ds.add_havea_footer(slide, page_num, total)


# ─── Slide 16 — Prochaines étapes ───────────────────────────────────────

def add_slide_next_steps(prs, page_num, total):
    slide = blank_slide(prs)
    ds.add_nav_breadcrumb(slide, NAV_SECTIONS, active_idx=6)
    ds.add_editorial_title(
        slide,
        "Prochaines étapes · roadmap 3 phases + 8 actions immédiates",
        size=18,
    )

    # Ligne 1 : Roadmap 3 phases
    phases = [
        ("PHASE 1 · MVP", "0–3 mois", "6 SKUs · toggle + page dédiée · engagement zéro · objectif 50-100 abonnés actifs, churn baseline à 90j"),
        ("PHASE 2 · OPT", "3–6 mois", "Ajout Whey (lock goût) · A/B wording · programme fidélité cumulable · analyse cohortes par SKU"),
        ("PHASE 3 · SCALE", "6–12 mois", "Tous les SKUs éligibles · parrainage abo · quiz d'onboarding · intégration programme ambassadeurs"),
    ]
    card_w = 3.1
    gap = 0.1
    start_x = 0.4
    for i, (title, dur, body) in enumerate(phases):
        x = start_x + i * (card_w + gap)
        ds.add_rect(slide, Inches(x), Inches(1.5), Inches(card_w), Inches(2.0), fill_color=ds.GREY_100)
        ds.add_rect(slide, Inches(x), Inches(1.5), Inches(0.08), Inches(2.0), fill_color=ds.MINT)
        ds.add_textbox(
            slide,
            title,
            Inches(x + 0.2),
            Inches(1.65),
            Inches(card_w - 0.4),
            Inches(0.25),
            font=ds.BODY_FONT,
            size=9,
            bold=True,
            color=ds.MINT,
        )
        ds.add_textbox(
            slide,
            dur,
            Inches(x + 0.2),
            Inches(1.9),
            Inches(card_w - 0.4),
            Inches(0.3),
            font=ds.DISPLAY_FONT,
            size=16,
            color=ds.INK,
        )
        ds.add_textbox(
            slide,
            body,
            Inches(x + 0.2),
            Inches(2.3),
            Inches(card_w - 0.4),
            Inches(1.15),
            font=ds.BODY_FONT,
            size=9,
            color=ds.INK,
            line_spacing=1.3,
        )

    # Bottom : 8 actions immédiates
    ds.add_textbox(
        slide,
        "8 ACTIONS IMMÉDIATES",
        Inches(0.4),
        Inches(3.7),
        Inches(9),
        Inches(0.25),
        font=ds.BODY_FONT,
        size=9,
        bold=True,
        color=ds.MINT,
    )
    actions = [
        "1. Validation comité projet du modèle commercial",
        "2. Audit Loi Chatel avec juriste (1-2 jours)",
        "3. Choix tech : Recharge / Shopify Native / Bold (devis)",
        "4. Export Shopify 90 jours · baseline potentiel abonnés",
        "5. Design / maquettes : fiche + page abo + email 48h",
        "6. Mise en place technique (2-3 semaines)",
        "7. Lancement soft · 6 SKUs · mesure 1 mois",
        "8. Lancement officiel · après ajustements UX",
    ]
    # 2 colonnes
    left_col = actions[:4]
    right_col = actions[4:]
    ds.add_bullets(slide, left_col, 0.4, 4.0, 4.6, 1.3, size=9, marker="·", marker_color=ds.MINT)
    ds.add_bullets(slide, right_col, 5.2, 4.0, 4.6, 1.3, size=9, marker="·", marker_color=ds.MINT)

    ds.add_havea_footer(slide, page_num, total)


# ─── Slide 17 (optionnel) — Closing ─────────────────────────────────────

def add_slide_closing(prs, page_num, total):
    slide = blank_slide(prs)
    # Fond sombre full
    ds.add_rect(slide, Inches(0), Inches(0), ds.SLIDE_W, ds.SLIDE_H, fill_color=ds.INK)
    # Feuille mint centrée
    ds.add_leaf_logo(slide, 4.7, 1.2, 1.2)
    # Titre
    ds.add_textbox(
        slide,
        "Merci",
        Inches(0),
        Inches(2.8),
        Inches(10),
        Inches(0.9),
        font=ds.DISPLAY_FONT,
        size=52,
        color=ds.WHITE,
        align=PP_ALIGN.CENTER,
        bold=False,
    )
    ds.add_rect(slide, Inches(4.55), Inches(3.75), Inches(0.9), Inches(0.04), fill_color=ds.MINT)
    ds.add_textbox(
        slide,
        "Benchmark abonnement · mission interne comité projet",
        Inches(0),
        Inches(3.9),
        Inches(10),
        Inches(0.3),
        font=ds.BODY_FONT,
        size=11,
        italic=True,
        color=ds.GREY_300,
        align=PP_ALIGN.CENTER,
    )
    ds.add_textbox(
        slide,
        "2026-04-15 · Impulse Nutrition · Havea Group",
        Inches(0),
        Inches(4.25),
        Inches(10),
        Inches(0.3),
        font=ds.BODY_FONT,
        size=9,
        color=ds.MINT,
        align=PP_ALIGN.CENTER,
    )
    # Havea logo
    ds.add_havea_logo_block(slide, x=4.35, y=4.7, w_inches=1.3, h_inches=0.6)


# ─── Main ───────────────────────────────────────────────────────────────

def load_brands() -> dict:
    brands = {}
    for slug in BRAND_ORDER:
        path = MARQUES / slug / "data.json"
        brands[slug] = json.loads(path.read_text(encoding="utf-8"))
    return brands


def main() -> None:
    brands = load_brands()

    prs = Presentation()
    prs.slide_width = ds.SLIDE_W
    prs.slide_height = ds.SLIDE_H

    total = 17  # fixed for footer display

    # 01 — Cover
    add_slide_cover(prs)

    # 02 — Methodology
    add_slide_methodology(prs, 2, total)

    # 03-10 — One slide per brand (8 slides)
    for i, slug in enumerate(BRAND_ORDER):
        add_slide_brand(prs, brands[slug], page_num=3 + i, total=total)

    # 11 — Synthesis table
    add_slide_synthesis(prs, brands, 11, total)

    # 12 — Profitability
    add_slide_profitability(prs, 12, total)

    # 13-15 — Recommandation (3 slides)
    add_slide_reco_products(prs, 13, total)
    add_slide_reco_commercial(prs, 14, total)
    add_slide_reco_ux(prs, 15, total)

    # 16 — Next steps
    add_slide_next_steps(prs, 16, total)

    # 17 — Closing
    add_slide_closing(prs, 17, total)

    prs.save(str(OUT))
    size_kb = OUT.stat().st_size // 1024
    print(f"\n[ok] Deck généré : {OUT.relative_to(ROOT.parent)} ({size_kb} KB · {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
