"""
deck_style.py — Design system + helpers réutilisables pour build_deck.py.

Reprend le style du deck Havea (palette mint, Source Sans Pro, photo hero,
layouts denses) via python-pptx. Toutes les fonctions prennent une `slide`
et des coordonnées en EMU (Inches convertis) et ajoutent des shapes.

Constantes principales :
  - Couleurs Havea (MINT, TEAL_DARK, INK, etc.)
  - Fonts (DISPLAY, BODY) avec fallback automatique Calibri
  - Dimensions slide (SLIDE_W, SLIDE_H) au format 10x5.625"

Helpers :
  - add_nav_breadcrumb(slide, sections, active_idx)
  - add_editorial_title(slide, title, left=0.4, top=0.5, width=9.2)
  - add_leaf_logo(slide, x, y, height)
  - add_havea_footer(slide, slide_number)
  - add_scorecard(slide, x, y, scores)
  - add_cover_hero(slide, photo_path, title, subtitle, week_label)
  - add_data_table(slide, headers, rows, x, y, w, col_widths=None)
  - add_textbox(slide, text, x, y, w, h, font_size, bold, color, align)
  - add_rect(slide, x, y, w, h, fill_color, line_color=None)
  - fit_image(image_path, box_w, box_h)  — utility for sizing
"""
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ─── Dimensions slide (16:9 Havea) ─────────────────────────────────────
SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.625)

# ─── Palette Havea ─────────────────────────────────────────────────────
MINT = RGBColor(0x27, 0x99, 0x89)          # Primary Havea mint
TEAL_DARK = RGBColor(0x01, 0x4F, 0x4F)     # Titres
TEAL_DEEP = RGBColor(0x00, 0x55, 0x55)     # Fonds blocs insight
GREEN_VERY_DARK = RGBColor(0x07, 0x33, 0x1D)
TEAL_CLEAR = RGBColor(0x00, 0x9C, 0xA0)
BLUE_SOFT = RGBColor(0x3E, 0x6F, 0x92)
BLUE_LIGHT = RGBColor(0x67, 0xA7, 0xC0)
BEIGE_GOLD = RGBColor(0xFE, 0xD9, 0x92)
SALMON = RGBColor(0xF0, 0x82, 0x76)
INK = RGBColor(0x0A, 0x0F, 0x1E)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY_900 = RGBColor(0x1E, 0x24, 0x33)
GREY_700 = RGBColor(0x4A, 0x52, 0x5E)
GREY_500 = RGBColor(0x6B, 0x65, 0x58)
GREY_300 = RGBColor(0xC7, 0xBF, 0xA9)
GREY_100 = RGBColor(0xF5, 0xF1, 0xE8)
CREAM = RGBColor(0xFA, 0xF7, 0xEF)

# Accent des scores /5
def score_color(score):
    if score is None:
        return GREY_500
    if score >= 4.0:
        return MINT
    if score >= 3.0:
        return BLUE_SOFT
    return SALMON


# ─── Fonts (Source Sans Pro avec fallback Calibri) ─────────────────────
DISPLAY_FONT = "Source Sans Pro Light"
DISPLAY_FONT_BOLD = "Source Sans Pro"
BODY_FONT = "Source Sans Pro"
MONO_FONT = "Menlo"  # pour chiffres/data

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "deck" / "assets"

# ─── Primitives ────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    """Ajoute un rectangle plein (utilitaire pour overlays, encadrés, fonds)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.shadow.inherit = False
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    return shape


def add_textbox(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    font=BODY_FONT,
    size=10,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    italic=False,
    line_spacing=1.1,
):
    """Ajoute un textbox stylé. Retourne le shape pour chaînage."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    para = tf.paragraphs[0]
    para.alignment = align
    para.line_spacing = line_spacing
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_multiline(slide, lines_spec, x, y, w, h, *, align=PP_ALIGN.LEFT, line_spacing=1.2):
    """lines_spec = list of (text, size, bold, color, italic) tuples."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, spec in enumerate(lines_spec):
        text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        italic = spec[4] if len(spec) > 4 else False
        font_name = spec[5] if len(spec) > 5 else BODY_FONT
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


# ─── Composants de slide récurrents ────────────────────────────────────

def add_havea_footer(slide, slide_number=None, total=None):
    """Footer bas de slide : bande mint fine + mention confidentialité + numéro."""
    # Bande fine mint en bas
    add_rect(slide, Inches(0), Inches(5.55), SLIDE_W, Inches(0.075), fill_color=MINT)
    # Mention confidentielle à gauche
    add_textbox(
        slide,
        "Internal document — Strictly confidential",
        Inches(0.4),
        Inches(5.37),
        Inches(5),
        Inches(0.2),
        font=BODY_FONT,
        size=7,
        italic=True,
        color=GREY_500,
    )
    # Numéro de slide à droite
    if slide_number is not None:
        label = f"{slide_number:02d}"
        if total:
            label += f" / {total:02d}"
        add_textbox(
            slide,
            label,
            Inches(9.3),
            Inches(5.37),
            Inches(0.6),
            Inches(0.2),
            font=BODY_FONT,
            size=8,
            bold=True,
            color=GREY_500,
            align=PP_ALIGN.RIGHT,
        )


def add_nav_breadcrumb(slide, sections, active_idx, *, top=0.15):
    """Nav top Havea : 'Section A  Section B  Section C' avec active en mint bold."""
    x = Inches(0.4)
    for i, name in enumerate(sections):
        is_active = i == active_idx
        color = MINT if is_active else GREY_500
        bold = is_active
        width = Inches(max(1.0, len(name) * 0.09))
        box = add_textbox(
            slide,
            name,
            x,
            Inches(top),
            width,
            Inches(0.2),
            font=BODY_FONT,
            size=8,
            bold=bold,
            color=color,
        )
        x = Emu(box.left + box.width + Inches(0.2))


def add_editorial_title(slide, title, *, left=0.4, top=0.45, width=9.2, size=20, color=INK):
    """Titre éditorial long type Havea — Source Sans Pro Light, gros, pas bold."""
    add_textbox(
        slide,
        title,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.55),
        font=DISPLAY_FONT,
        size=size,
        bold=False,
        color=color,
        line_spacing=1.1,
    )
    # Underline mint fine sous le titre
    add_rect(
        slide,
        Inches(left),
        Inches(top + 0.75),
        Inches(0.6),
        Inches(0.04),
        fill_color=MINT,
    )


def add_leaf_logo(slide, x, y, height_inches):
    """Ajoute le logo feuille mint Havea. x/y en inches."""
    path = str(ASSETS / "leaf_mint.png")
    return slide.shapes.add_picture(path, Inches(x), Inches(y), height=Inches(height_inches))


def add_havea_logo_block(slide, x, y, w_inches=1.3, h_inches=0.6):
    """Bloc 'HAVEA' stylisé : rectangle mint + texte blanc (recrée le logo)."""
    add_rect(slide, Inches(x), Inches(y), Inches(w_inches), Inches(h_inches), fill_color=MINT)
    add_textbox(
        slide,
        "Havea",
        Inches(x),
        Inches(y + 0.1),
        Inches(w_inches),
        Inches(0.3),
        font=DISPLAY_FONT,
        size=16,
        bold=False,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        "— G R O U P —",
        Inches(x),
        Inches(y + 0.38),
        Inches(w_inches),
        Inches(0.18),
        font=BODY_FONT,
        size=7,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )


# ─── Cover slide layout ────────────────────────────────────────────────

def add_cover_hero(slide, photo_path, title, subtitle, week_label):
    """Layout complet cover : photo hero fond + overlay noir + panneau info + logos."""
    # 1. Photo hero plein écran
    if photo_path and Path(photo_path).exists():
        slide.shapes.add_picture(str(photo_path), Inches(0), Inches(0), width=SLIDE_W, height=SLIDE_H)
    # 2. Overlay noir semi-transparent sur la partie droite (panneau info)
    # python-pptx ne supporte pas la transparence directe sur shapes, on utilise un gros rectangle
    # semi-dark — workaround: rectangle noir plein sur panneau info + texte blanc
    panel_x = Inches(4.5)
    panel_w = Inches(5.5)
    add_rect(slide, panel_x, Inches(0), panel_w, SLIDE_H, fill_color=INK)

    # 3. Logo feuille mint en haut du panneau
    add_leaf_logo(slide, 5.0, 0.6, 1.0)

    # 4. Titre principal (IMPULSE NUTRITION style)
    add_textbox(
        slide,
        title,
        Inches(4.7),
        Inches(2.0),
        Inches(5.1),
        Inches(0.9),
        font=DISPLAY_FONT,
        size=34,
        bold=False,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        line_spacing=1.0,
    )
    # Séparateur mint
    add_rect(slide, Inches(6.2), Inches(2.95), Inches(2.1), Inches(0.04), fill_color=MINT)

    # 5. Sous-titre
    add_textbox(
        slide,
        subtitle,
        Inches(4.7),
        Inches(3.1),
        Inches(5.1),
        Inches(0.4),
        font=BODY_FONT,
        size=13,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        italic=True,
    )

    # 6. Week label
    if week_label:
        add_textbox(
            slide,
            week_label,
            Inches(4.7),
            Inches(3.65),
            Inches(5.1),
            Inches(0.3),
            font=BODY_FONT,
            size=11,
            color=MINT,
            align=PP_ALIGN.CENTER,
            bold=True,
        )

    # 7. Mention confidentialité (bottom-left sur photo)
    add_textbox(
        slide,
        "Internal document — Strictly confidential",
        Inches(0.2),
        Inches(5.37),
        Inches(4),
        Inches(0.2),
        font=BODY_FONT,
        size=7,
        italic=True,
        color=WHITE,
    )

    # 8. Logo Havea bottom-right sur fond mint
    add_havea_logo_block(slide, x=8.35, y=4.85, w_inches=1.3, h_inches=0.6)


# ─── Scorecard /5 ──────────────────────────────────────────────────────

def add_scorecard(slide, x, y, scores, *, cell_w=1.3, cell_h=0.85, gap=0.1):
    """Scorecard 4 cellules : UX, Offre, Pertinence, Global.
    scores = dict {ux, offre, pertinence, global} avec valeurs float or None."""
    labels = [
        ("UX", scores.get("ux")),
        ("Offre", scores.get("offre")),
        ("Pertinence", scores.get("pertinence")),
        ("Global", scores.get("global")),
    ]
    for i, (label, val) in enumerate(labels):
        cx = Inches(x + i * (cell_w + gap))
        cy = Inches(y)
        is_global = label == "Global"
        bg = MINT if is_global else GREY_100
        label_color = WHITE if is_global else GREY_700
        value_color = WHITE if is_global else score_color(val)

        add_rect(slide, cx, cy, Inches(cell_w), Inches(cell_h), fill_color=bg)
        # Label
        add_textbox(
            slide,
            label.upper(),
            cx,
            Emu(cy + Inches(0.1)),
            Inches(cell_w),
            Inches(0.2),
            font=BODY_FONT,
            size=8,
            bold=True,
            color=label_color,
            align=PP_ALIGN.CENTER,
        )
        # Value
        val_text = f"{val}" if val is not None else "—"
        add_textbox(
            slide,
            val_text,
            cx,
            Emu(cy + Inches(0.3)),
            Inches(cell_w),
            Inches(0.5),
            font=DISPLAY_FONT,
            size=26,
            bold=False,
            color=value_color,
            align=PP_ALIGN.CENTER,
        )
        # /5 suffix
        if val is not None:
            add_textbox(
                slide,
                "/ 5",
                cx,
                Emu(cy + Inches(0.58)),
                Inches(cell_w),
                Inches(0.2),
                font=BODY_FONT,
                size=7,
                color=label_color,
                align=PP_ALIGN.CENTER,
            )


# ─── Data table (dense, style Havea) ───────────────────────────────────

def add_data_table(slide, headers, rows, x, y, w, *, col_widths=None, row_h=0.3, header_h=0.3, font_size=8):
    """Tableau compact sans bordures intérieures lourdes, style Havea."""
    n_cols = len(headers)
    total_w = w
    if col_widths is None:
        col_widths = [total_w / n_cols] * n_cols

    # Header row (underline mint)
    cx = x
    for i, h in enumerate(headers):
        add_textbox(
            slide,
            h.upper(),
            Inches(cx),
            Inches(y),
            Inches(col_widths[i]),
            Inches(header_h),
            font=BODY_FONT,
            size=font_size - 1,
            bold=True,
            color=GREY_700,
        )
        cx += col_widths[i]
    # Underline header
    add_rect(slide, Inches(x), Inches(y + header_h + 0.02), Inches(w), Inches(0.02), fill_color=INK)

    # Data rows
    current_y = y + header_h + 0.08
    for row in rows:
        cx = x
        for i, cell in enumerate(row):
            text = str(cell) if cell is not None else "—"
            add_textbox(
                slide,
                text,
                Inches(cx),
                Inches(current_y),
                Inches(col_widths[i]),
                Inches(row_h),
                font=BODY_FONT,
                size=font_size,
                color=INK,
            )
            cx += col_widths[i]
        # Fine rule under row
        add_rect(
            slide,
            Inches(x),
            Inches(current_y + row_h + 0.01),
            Inches(w),
            Inches(0.005),
            fill_color=GREY_300,
        )
        current_y += row_h + 0.05
    return current_y  # return y position after table


# ─── Screenshot + caption ──────────────────────────────────────────────

def add_screenshot_with_caption(slide, image_path, x, y, w, h, *, caption=None):
    """Ajoute un screenshot centré dans un cadre avec légende optionnelle en dessous."""
    if not image_path or not Path(image_path).exists():
        # Placeholder si manquant
        add_rect(slide, Inches(x), Inches(y), Inches(w), Inches(h), fill_color=GREY_100, line_color=GREY_300)
        add_textbox(
            slide,
            "[screenshot manquant]",
            Inches(x),
            Inches(y + h / 2),
            Inches(w),
            Inches(0.3),
            font=BODY_FONT,
            size=10,
            italic=True,
            color=GREY_500,
            align=PP_ALIGN.CENTER,
        )
        return

    # Cadre ink fine
    add_rect(slide, Inches(x), Inches(y), Inches(w), Inches(h), line_color=GREY_300, line_width=Emu(6000))
    # Ajout de l'image avec ratio préservé, inscrite dans le cadre (laisse 0.1 inch de marge interne)
    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(x + 0.05),
        Inches(y + 0.05),
        width=Inches(w - 0.1),
    )
    # Si l'image est trop haute, on re-size par height
    if pic.height > Inches(h - 0.1):
        # Remove and re-add with height constraint
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            str(image_path),
            Inches(x + 0.05),
            Inches(y + 0.05),
            height=Inches(h - 0.1),
        )
        # Recenter horizontally
        new_left = Inches(x) + (Inches(w) - pic.width) // 2
        pic.left = new_left
    else:
        # Recenter vertically if image is smaller than box
        new_top = Inches(y) + (Inches(h) - pic.height) // 2
        pic.top = new_top

    # Caption en dessous
    if caption:
        add_textbox(
            slide,
            caption,
            Inches(x),
            Inches(y + h + 0.05),
            Inches(w),
            Inches(0.25),
            font=BODY_FONT,
            size=8,
            italic=True,
            color=GREY_500,
            align=PP_ALIGN.CENTER,
        )


# ─── Bullets list ──────────────────────────────────────────────────────

def add_bullets(slide, items, x, y, w, h, *, size=10, color=INK, marker="•", marker_color=None):
    """Liste à puces stylée. items = list of strings."""
    if marker_color is None:
        marker_color = MINT
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.line_spacing = 1.2
        # Marker
        run_m = p.add_run()
        run_m.text = f"{marker}  "
        run_m.font.name = BODY_FONT
        run_m.font.size = Pt(size)
        run_m.font.bold = True
        run_m.font.color.rgb = marker_color
        # Text
        run_t = p.add_run()
        run_t.text = item
        run_t.font.name = BODY_FONT
        run_t.font.size = Pt(size)
        run_t.font.color.rgb = color
    return box


# ─── Insight callout (bloc sombre + accent mint) ───────────────────────

def add_insight_box(slide, x, y, w, h, title, body, *, bg=INK, accent=MINT):
    """Bloc insight type Havea : fond sombre + accent mint à gauche + titre + corps."""
    # Fond
    add_rect(slide, Inches(x), Inches(y), Inches(w), Inches(h), fill_color=bg)
    # Accent mint à gauche
    add_rect(slide, Inches(x), Inches(y), Inches(0.08), Inches(h), fill_color=accent)
    # Titre (kicker)
    add_textbox(
        slide,
        title.upper(),
        Inches(x + 0.25),
        Inches(y + 0.2),
        Inches(w - 0.5),
        Inches(0.25),
        font=BODY_FONT,
        size=8,
        bold=True,
        color=accent,
    )
    # Body
    add_textbox(
        slide,
        body,
        Inches(x + 0.25),
        Inches(y + 0.5),
        Inches(w - 0.5),
        Inches(h - 0.7),
        font=BODY_FONT,
        size=11,
        color=WHITE,
        line_spacing=1.3,
    )
