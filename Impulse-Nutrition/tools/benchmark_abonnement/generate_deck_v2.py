#!/usr/bin/env python3
"""
Benchmark Abonnement Nutrition FR 2026 — Deck V2
Generates a comprehensive PowerPoint from master_data.json + ux_deep.json + rentabilite_output.json
"""

import json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

BASE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(BASE)

# Brand colors
MARINE = RGBColor(0x1B, 0x2B, 0x4B)
ORANGE = RGBColor(0xF0, 0x5A, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
LIGHT_BLUE = RGBColor(0xEB, 0xF0, 0xF7)

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def load_data():
    with open(os.path.join(BASE, "master_data.json"), encoding="utf-8") as f:
        master = json.load(f)
    with open(os.path.join(BASE, "rentabilite_output.json"), encoding="utf-8") as f:
        renta = json.load(f)
    return master, renta


def add_background(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=12, color=DARK_GRAY,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=11, color=DARK_GRAY, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2), bullet=False):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    if bullet:
        p.level = 0
    return p


def add_header_bar(slide, title_text, subtitle_text=""):
    """Add marine header bar with title."""
    bar = add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=MARINE)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.55),
                title_text, font_size=26, color=WHITE, bold=True)
    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.65), Inches(10), Inches(0.35),
                    subtitle_text, font_size=13, color=RGBColor(0xBB, 0xCC, 0xDD))
    # Orange accent line
    add_shape(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), fill_color=ORANGE)


def add_footer(slide, page_num, total_pages):
    """Add footer with page number and branding."""
    add_shape(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), fill_color=LIGHT_GRAY)
    add_textbox(slide, Inches(0.6), Inches(7.18), Inches(4), Inches(0.28),
                "Impulse Nutrition — Benchmark Abonnement 2026", font_size=8, color=MID_GRAY)
    add_textbox(slide, Inches(11), Inches(7.18), Inches(2), Inches(0.28),
                f"{page_num} / {total_pages}", font_size=8, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


def add_kpi_card(slide, left, top, width, height, label, value, accent_color=ORANGE):
    """Add a KPI card with colored accent."""
    card = add_shape(slide, left, top, width, height, fill_color=WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD), border_width=Pt(1))
    # Accent top bar
    add_shape(slide, left, top, width, Inches(0.06), fill_color=accent_color)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.12), width - Inches(0.24), Inches(0.35),
                value, font_size=22, color=MARINE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.48), width - Inches(0.24), Inches(0.3),
                label, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def build_table(slide, left, top, width, rows_data, col_widths, header_color=MARINE):
    """Build a styled table."""
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(rows * 0.35))
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.name = "Calibri"
                if ri == 0:
                    p.font.color.rgb = WHITE
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.font.color.rgb = DARK_GRAY
                    p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE

    return table_shape


# ============= SLIDE BUILDERS =============

def slide_title(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide, MARINE)

    # Big title
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                "Benchmark Abonnement", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.8),
                "Nutrition Sportive — Marché Français 2026", font_size=28, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    # Orange divider
    add_shape(slide, Inches(5), Inches(3.6), Inches(3.3), Inches(0.04), fill_color=ORANGE)

    # Subtitle
    add_textbox(slide, Inches(2), Inches(3.9), Inches(9), Inches(0.6),
                "Analyse comparative de 8 marques · Modèles commerciaux · UX · Logistique · Rentabilité",
                font_size=14, color=RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)

    # Bottom info
    add_textbox(slide, Inches(2), Inches(5.2), Inches(9), Inches(0.4),
                "HAVEA Commercial Services — Impulse Nutrition", font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2), Inches(5.6), Inches(9), Inches(0.4),
                "Avril 2026 · Usage interne — Comité Projet", font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Logo
    logo_path = os.path.join(ROOT, "assets", "logo_impulse.jpeg")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(5.4), Inches(6.1), Inches(2.5))


def slide_context(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_header_bar(slide, "Contexte & Objectif")
    add_footer(slide, 2, total)

    # Context
    tb1 = add_textbox(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.5),
                      "", font_size=12, color=DARK_GRAY)
    tf = tb1.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Contexte"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MARINE

    texts = [
        "Impulse Nutrition envisage de lancer un modèle d'abonnement sur impulse-nutrition.fr.",
        "L'abonnement est un levier de croissance majeur dans la nutrition sportive : LTV x3-5 vs achat ponctuel.",
        "Le marché FR de la nutrition sportive spécialisée pèse ~210 M€ (2023) avec un CAGR de 8.3%.",
        "6 marques sur 8 analysées proposent déjà un abonnement — Impulse doit se positionner."
    ]
    for t in texts:
        add_paragraph(tf, f"• {t}", font_size=11, color=DARK_GRAY, space_before=Pt(6))

    # Objective
    tb2 = add_textbox(slide, Inches(7), Inches(1.4), Inches(5.8), Inches(2.5),
                      "", font_size=12, color=DARK_GRAY)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "Objectifs de cette étude"
    tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = MARINE

    objs = [
        "Benchmarker les modèles d'abonnement existants (8 marques FR/internationales)",
        "Analyser l'UX, le wording, le placement du module abonnement",
        "Évaluer la rentabilité (coûts logistiques vs remise)",
        "Formuler une recommandation concrète et opérationnelle pour Impulse"
    ]
    for o in objs:
        add_paragraph(tf2, f"• {o}", font_size=11, color=DARK_GRAY, space_before=Pt(6))

    # KPI cards
    add_kpi_card(slide, Inches(0.6), Inches(4.6), Inches(2.8), Inches(0.85), "Marques analysées", "8")
    add_kpi_card(slide, Inches(3.8), Inches(4.6), Inches(2.8), Inches(0.85), "Ont un abonnement", "6 / 8", accent_color=GREEN)
    add_kpi_card(slide, Inches(7), Inches(4.6), Inches(2.8), Inches(0.85), "LTV abonné vs one-shot", "x3 à x5")
    add_kpi_card(slide, Inches(10.2), Inches(4.6), Inches(2.8), Inches(0.85), "Marché FR spécialisé", "210 M€")

    # Bottom box
    box = add_shape(slide, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.1),
                    fill_color=LIGHT_BLUE, border_color=MARINE, border_width=Pt(1))
    add_textbox(slide, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.9),
                "💡 Enjeu central : à partir de quel ticket moyen l'abonnement avec livraison offerte est-il rentable ? "
                "→ Analyse de rentabilité complète incluse dans ce deck (seuil estimé : 25€ minimum).",
                font_size=11, color=MARINE)


def slide_methodology(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_header_bar(slide, "Méthodologie & Périmètre", "8 marques · 7 axes d'analyse · Sources web publiques")
    add_footer(slide, 3, total)

    # Brands grid
    brands_info = [
        ("Nutri&Co", "✅ Abonnement", "Shopify", "Ref. FR premium"),
        ("Nutrimuscle", "✅ Abonnement", "Shopify Plus", "Leader volume"),
        ("Novoma", "✅ Abonnement", "Shopify", "Santé-sport"),
        ("Nutripure", "❌ Absent", "PrestaShop", "Concurrent direct"),
        ("Cuure", "✅ Box perso", "Custom", "Santé premium"),
        ("Decathlon", "✅ Abonnement", "Custom", "Grande distribution"),
        ("Aqeelab", "✅ Abonnement", "Shopify", "Pure player DTC"),
        ("MyProtein", "✅ Subscribe&Save", "Custom (THG)", "International"),
    ]

    rows = [["Marque", "Abonnement", "Plateforme", "Positionnement"]] + [list(b) for b in brands_info]
    build_table(slide, Inches(0.6), Inches(1.4), Inches(7), rows,
                [Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.9)])

    # Methodology axes
    tb = add_textbox(slide, Inches(8.2), Inches(1.4), Inches(4.8), Inches(5),
                     "", font_size=12)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "7 axes d'analyse"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MARINE

    axes = [
        ("1. Catalogues & pricing", "Prix produit par produit, deltas abo vs one-shot"),
        ("2. Logistique & transport", "Grilles transporteurs FR, seuils livraison gratuite"),
        ("3. Tech stack", "Plateformes e-commerce, apps abonnement"),
        ("4. Avis clients", "Trustpilot scores, verbatims, NPS"),
        ("5. Marché & tendances", "Taille marché, churn, LTV, comportement consommateur"),
        ("6. UX & wording", "Screenshots, CTA, placement module, pages dédiées"),
        ("7. Rentabilité", "Modèle de marge, seuils de rentabilité"),
    ]
    for title, desc in axes:
        add_paragraph(tf, title, font_size=11, color=ORANGE, bold=True, space_before=Pt(8))
        add_paragraph(tf, desc, font_size=9, color=MID_GRAY, space_before=Pt(1))


def slide_brand(prs, brand_key, brand_data, page_num, total):
    """Build a comprehensive slide for one brand."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    cat = brand_data.get("catalogue", {})
    log = brand_data.get("logistique", {})
    avis = brand_data.get("avis", {})
    ux = brand_data.get("ux_deep", {})
    tech = brand_data.get("tech", {})

    brand_name = cat.get("name", brand_key)
    has_sub = cat.get("has_subscription", False)

    sub_status = "✅ Abonnement actif" if has_sub else "❌ Pas d'abonnement"
    add_header_bar(slide, brand_name, f"{sub_status} · {tech.get('platform', 'N/A')} · Trustpilot {avis.get('trustpilot_score', 'N/A')}/5")
    add_footer(slide, page_num, total)

    # LEFT COLUMN: Key metrics
    y = Inches(1.35)

    # Subscription model card
    card = add_shape(slide, Inches(0.4), y, Inches(4.2), Inches(1.6),
                     fill_color=LIGHT_BLUE, border_color=MARINE, border_width=Pt(1))
    tb_model = add_textbox(slide, Inches(0.55), y + Inches(0.08), Inches(3.9), Inches(1.4), "")
    tf = tb_model.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Modèle commercial"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MARINE

    if has_sub:
        discount = cat.get("subscription_discount_pct", "N/A")
        discount_str = f"{discount}%" if isinstance(discount, (int, float)) else str(discount)
        model_info = [
            f"Remise : {discount_str}",
            f"Fréquences : {', '.join(cat.get('frequencies_available', ['N/A']))}",
            f"Engagement : {cat.get('flexibility', {}).get('minimum_commitment', 'N/A')}",
            f"Pause/Modif : {'✅' if cat.get('flexibility', {}).get('pause') else '❌'} / {'✅' if cat.get('flexibility', {}).get('modify') else '❌'}"
        ]
    else:
        model_info = [
            cat.get("subscription_model", "Aucun abonnement"),
            cat.get("why_no_subscription", "")[:100]
        ]

    for info in model_info:
        add_paragraph(tf, f"  {info}", font_size=10, color=DARK_GRAY, space_before=Pt(3))

    # Shipping card
    y2 = Inches(3.1)
    card2 = add_shape(slide, Inches(0.4), y2, Inches(4.2), Inches(1.1),
                      fill_color=WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD), border_width=Pt(1))
    tb_ship = add_textbox(slide, Inches(0.55), y2 + Inches(0.08), Inches(3.9), Inches(0.9), "")
    tf2 = tb_ship.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "Livraison"
    tf2.paragraphs[0].font.size = Pt(13)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = MARINE

    relais = log.get("frais_relais", log.get("frais_standard", "N/A"))
    seuil = log.get("seuil_gratuit_relais", log.get("seuil_gratuit", "N/A"))
    add_paragraph(tf2, f"  Point relais : {relais}€ · Gratuit dès {seuil}€", font_size=10, color=DARK_GRAY, space_before=Pt(3))

    domicile = log.get("frais_domicile", log.get("frais_domicile_standard", "N/A"))
    seuil_dom = log.get("seuil_gratuit_domicile", "N/A")
    add_paragraph(tf2, f"  Domicile : {domicile}€ · Gratuit dès {seuil_dom}€", font_size=10, color=DARK_GRAY, space_before=Pt(2))

    # UX & Wording card
    y3 = Inches(4.35)
    sub_ux = ux.get("subscription_ux", {})
    card3 = add_shape(slide, Inches(0.4), y3, Inches(4.2), Inches(1.5),
                      fill_color=WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD), border_width=Pt(1))
    tb_ux = add_textbox(slide, Inches(0.55), y3 + Inches(0.08), Inches(3.9), Inches(1.3), "")
    tf3 = tb_ux.text_frame
    tf3.word_wrap = True
    tf3.paragraphs[0].text = "UX & Wording"
    tf3.paragraphs[0].font.size = Pt(13)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = MARINE

    wording = sub_ux.get("wording_pitch", sub_ux.get("wording_cta", "N/A"))
    add_paragraph(tf3, f'  CTA : "{sub_ux.get("wording_cta", "N/A")}"', font_size=10, color=DARK_GRAY, space_before=Pt(3))
    add_paragraph(tf3, f'  Pitch : "{wording[:60]}"', font_size=10, color=DARK_GRAY, space_before=Pt(2))
    placement = sub_ux.get("placement_on_product_page", "N/A")
    add_paragraph(tf3, f"  Placement : {placement[:55]}", font_size=9, color=MID_GRAY, space_before=Pt(2))
    page_ded = "Oui" if sub_ux.get("dedicated_page") else "Non"
    add_paragraph(tf3, f"  Page dédiée : {page_ded}", font_size=10, color=DARK_GRAY, space_before=Pt(2))

    # RIGHT COLUMN: Screenshot + evaluation
    screenshot_dir = os.path.join(ROOT, "assets", "screenshots", "benchmark", "ux_deep")
    # Pick best screenshot (product_sub preferred)
    screenshots = ux.get("screenshots", [])
    screenshot_file = None
    for pref in ["product_sub", "abonnement_page", "quiz", "abo_homepage", "homepage"]:
        for s in screenshots:
            if pref in s:
                screenshot_file = s
                break
        if screenshot_file:
            break
    if not screenshot_file and screenshots:
        screenshot_file = screenshots[0]

    if screenshot_file:
        img_path = os.path.join(screenshot_dir, screenshot_file)
        if os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, Inches(4.9), Inches(1.35), Inches(4.5), Inches(3.3))
            except Exception:
                pass

    # Evaluation card (right side)
    y4 = Inches(1.35)
    ux_score = ux.get("overall_ux_score", "N/A")
    score_str = f"{ux_score}/5" if isinstance(ux_score, (int, float)) else "N/A"

    # KPI row at top right
    add_kpi_card(slide, Inches(9.7), y4, Inches(1.6), Inches(0.75), "UX Score", score_str, accent_color=ORANGE)
    tp_score = avis.get("trustpilot_score", "N/A")
    add_kpi_card(slide, Inches(11.5), y4, Inches(1.6), Inches(0.75), "Trustpilot", f"{tp_score}/5", accent_color=GREEN)

    # Strengths & weaknesses
    y5 = Inches(2.3)
    strengths = ux.get("strengths", [])
    weaknesses = ux.get("weaknesses", [])

    tb_eval = add_textbox(slide, Inches(9.6), y5, Inches(3.5), Inches(4.5), "")
    tfe = tb_eval.text_frame
    tfe.word_wrap = True
    tfe.paragraphs[0].text = "Points forts"
    tfe.paragraphs[0].font.size = Pt(12)
    tfe.paragraphs[0].font.bold = True
    tfe.paragraphs[0].font.color.rgb = GREEN

    for s in strengths[:4]:
        add_paragraph(tfe, f"✅ {s}", font_size=9, color=DARK_GRAY, space_before=Pt(3))

    add_paragraph(tfe, "", font_size=6)
    add_paragraph(tfe, "Points faibles", font_size=12, color=RED, bold=True, space_before=Pt(8))

    for w in weaknesses[:4]:
        add_paragraph(tfe, f"⚠️ {w}", font_size=9, color=DARK_GRAY, space_before=Pt(3))

    # Products sample
    products = cat.get("products", [])
    if products:
        y6 = Inches(6.0)
        prod_texts = []
        for p in products[:5]:
            name = p.get("name", "")
            price = p.get("price_oneshot", p.get("price_range", ""))
            price_sub = p.get("price_sub", "")
            if price_sub:
                prod_texts.append(f"{name}: {price}€ → {price_sub}€")
            elif price:
                prod_texts.append(f"{name}: {price}")

        add_textbox(slide, Inches(0.4), y6, Inches(12.8), Inches(0.3),
                    "Produits clés : " + " · ".join(prod_texts),
                    font_size=8, color=MID_GRAY)


def slide_synthesis_table(prs, master, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_header_bar(slide, "Tableau de synthèse comparatif", "Vue d'ensemble des 8 marques analysées")
    add_footer(slide, page_num, total)

    brand_order = ["nutriandco", "nutrimuscle", "novoma", "nutripure", "cuure", "decathlon", "aqeelab", "myprotein"]

    header = ["Marque", "Remise", "Fréquences", "Engagement", "Livr. gratuite", "Page dédiée", "Plateforme", "Trustpilot", "UX /5"]
    rows = [header]

    for bk in brand_order:
        bd = master["brands"].get(bk, {})
        cat = bd.get("catalogue", {})
        ux = bd.get("ux_deep", {})
        tech = bd.get("tech", {})
        avis = bd.get("avis", {})

        name = cat.get("name", bk)
        discount = cat.get("subscription_discount_pct", "❌")
        if isinstance(discount, (int, float)):
            discount = f"-{discount}%"
        elif discount == "❌":
            discount = "❌ Aucun"

        freqs = cat.get("frequencies_available", [])
        freq_str = f"{len(freqs)} ({freqs[0]}→{freqs[-1]})" if freqs else "N/A"

        engagement = cat.get("flexibility", {}).get("minimum_commitment", "N/A")
        livr = "Oui (abo)" if bk == "nutrimuscle" else ("Oui (relais)" if bk == "decathlon" else ("Incluse" if bk == "cuure" else "Seuil"))

        page_ded = "✅" if ux.get("subscription_ux", {}).get("dedicated_page") else "❌"
        platform = tech.get("platform", "N/A")
        tp = f"{avis.get('trustpilot_score', '?')}/5"

        ux_score = ux.get("overall_ux_score", "N/A")
        ux_str = f"{ux_score}" if isinstance(ux_score, (int, float)) else "N/A"

        rows.append([name, discount, freq_str, engagement, livr, page_ded, platform, tp, ux_str])

    col_widths = [Inches(1.4), Inches(1.2), Inches(1.5), Inches(1.2), Inches(1.3), Inches(1.0), Inches(1.3), Inches(1.0), Inches(0.8)]
    build_table(slide, Inches(0.4), Inches(1.4), Inches(12.5), rows, col_widths)

    # Key takeaways
    add_textbox(slide, Inches(0.4), Inches(5.2), Inches(12.5), Inches(1.5),
                "", font_size=11)
    tb = slide.shapes[-1]
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Points clés de la synthèse"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MARINE

    takeaways = [
        "La remise médiane est de 10-15% — Aqeelab (20%) et MyProtein (30-45%) sont les plus agressifs",
        "6/8 marques proposent un abonnement sans engagement — c'est le standard du marché",
        "Nutrimuscle se démarque avec la livraison gratuite incluse dans l'abonnement",
        "4/8 marques sont sur Shopify (3 via l'agence Axome) — écosystème dominant pour les DNVB nutrition",
        "Nutripure (PrestaShop) n'a pas d'abonnement : probable frein technique"
    ]
    for t in takeaways:
        add_paragraph(tf, f"→ {t}", font_size=10, color=DARK_GRAY, space_before=Pt(4))


def slide_ux_comparative(prs, master, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_header_bar(slide, "Analyse UX comparative", "Wording · Placement · Pages dédiées · Best practices")
    add_footer(slide, page_num, total)

    ux_comp = master.get("transversal", {}).get("ux_comparative", {})

    # UX ranking table
    ranking = ux_comp.get("ux_ranking", [])
    rows = [["Rang", "Marque", "Score /5", "Justification"]]
    for i, r in enumerate(ranking):
        rows.append([f"#{i+1}", r["brand"], str(r["score"]), r["reason"]])
    col_widths = [Inches(0.7), Inches(1.5), Inches(0.9), Inches(4.5)]
    build_table(slide, Inches(0.4), Inches(1.4), Inches(7.6), rows, col_widths)

    # Best practices column
    tb = add_textbox(slide, Inches(8.4), Inches(1.4), Inches(4.7), Inches(5.5), "")
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Best practices identifiées"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MARINE

    bp = [
        ("🏆 Meilleure page dédiée", ux_comp.get("best_dedicated_page", "")),
        ("🎯 Meilleur toggle produit", ux_comp.get("best_product_toggle", "")),
        ("🔄 Meilleur onboarding", ux_comp.get("best_onboarding", "")),
        ("💰 Remise la plus agressive", ux_comp.get("most_aggressive_discount", "")),
        ("📦 Meilleur incentive livraison", ux_comp.get("best_shipping_incentive", "")),
        ("📅 Meilleure flexibilité", ux_comp.get("best_frequency_flexibility", "")),
        ("⚠️ UX la plus faible", ux_comp.get("worst_ux", "")),
    ]
    for title, desc in bp:
        add_paragraph(tf, title, font_size=11, color=ORANGE, bold=True, space_before=Pt(8))
        add_paragraph(tf, desc[:85], font_size=9, color=DARK_GRAY, space_before=Pt(2))

    # Wording patterns
    patterns = ux_comp.get("wording_patterns", {})
    add_textbox(slide, Inches(0.4), Inches(5.0), Inches(7.6), Inches(2),
                "", font_size=11)
    tb2 = slide.shapes[-1]
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "Patterns de wording"
    tf2.paragraphs[0].font.size = Pt(13)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = MARINE

    add_paragraph(tf2, f"Wording dominant : \"{patterns.get('most_common', '')}\"", font_size=10, color=DARK_GRAY, space_before=Pt(4))
    add_paragraph(tf2, f"Argumentaire santé : {patterns.get('argumentaire_sante', '')}", font_size=10, color=DARK_GRAY, space_before=Pt(3))
    add_paragraph(tf2, f"Argumentaire prix : {patterns.get('argumentaire_prix', '')}", font_size=10, color=DARK_GRAY, space_before=Pt(3))
    add_paragraph(tf2, f"Argumentaire confort : {patterns.get('argumentaire_confort', '')}", font_size=10, color=DARK_GRAY, space_before=Pt(3))


def slide_logistics(prs, master, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_header_bar(slide, "Analyse logistique & livraison", "Coûts transporteurs · Seuils livraison gratuite · Avantages abonnés")
    add_footer(slide, page_num, total)

    # Carrier grid
    carriers_data = master.get("transversal", {}).get("carrier_grids_fr_2026", {})
    rows_carrier = [["Transporteur", "0-500g", "500g-1kg", "1-2kg", "2-5kg", "5-10kg"]]
    for carrier_name, grid in carriers_data.items():
        row = [carrier_name.replace("_", " ").title()]
        for weight_cat in ["0_500g", "500g_1kg", "1_2kg", "2_5kg", "5_10kg"]:
            val = grid.get(weight_cat, "N/A")
            row.append(f"{val}€" if isinstance(val, (int, float)) else str(val))
        rows_carrier.append(row)

    if len(rows_carrier) > 1:
        col_widths = [Inches(2), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2)]
        build_table(slide, Inches(0.4), Inches(1.4), Inches(8), rows_carrier, col_widths)

    # Brand shipping thresholds
    brand_order = ["nutriandco", "nutrimuscle", "novoma", "nutripure", "cuure", "decathlon", "aqeelab", "myprotein"]
    rows_ship = [["Marque", "Relais", "Seuil gratuit", "Domicile", "Seuil gratuit", "Avantage abo"]]
    for bk in brand_order:
        bd = master["brands"].get(bk, {})
        log = bd.get("logistique", {})
        cat = bd.get("catalogue", {})
        name = cat.get("name", bk)
        fr = log.get("frais_relais", log.get("frais_standard", "N/A"))
        sr = log.get("seuil_gratuit_relais", log.get("seuil_gratuit", "N/A"))
        fd = log.get("frais_domicile", log.get("frais_domicile_standard", "N/A"))
        sd = log.get("seuil_gratuit_domicile", "N/A")
        # Subscription shipping advantage
        adv_abo = log.get("avantage_abonnement", {})
        if adv_abo:
            adv = f"Dom. {adv_abo.get('frais_domicile_abo', '?')}€, gratuit dès {adv_abo.get('seuil_gratuit_domicile_abo', '?')}€"
        elif bk == "nutrimuscle":
            adv = "Livraison gratuite"
        elif bk == "decathlon":
            adv = "Gratuit en relais"
        elif bk == "cuure":
            adv = "Incluse dans abo"
        else:
            adv = "—"
        rows_ship.append([name, f"{fr}€", f"{sr}€", f"{fd}€", f"{sd}€", adv])

    col_widths2 = [Inches(1.4), Inches(0.9), Inches(1.1), Inches(0.9), Inches(1.1), Inches(2.8)]
    y_offset = Inches(1.4) if len(rows_carrier) <= 1 else Inches(4.2)
    build_table(slide, Inches(0.4), y_offset, Inches(12.5), rows_ship, col_widths2)

    # Key insight
    add_shape(slide, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.65),
              fill_color=LIGHT_BLUE, border_color=MARINE, border_width=Pt(1))
    add_textbox(slide, Inches(0.7), Inches(6.35), Inches(12), Inches(0.55),
                "💡 Novoma est le seul à proposer un avantage livraison différencié pour les abonnés (3.50€ vs 4.90€, gratuit dès 50€ vs 90€). "
                "Nutrimuscle offre la livraison 100% gratuite sur l'abo. → Levier puissant de conversion.",
                font_size=10, color=MARINE)


def slide_profitability(prs, renta, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_header_bar(slide, "Analyse de rentabilité", "Modèle de marge avec remise -10% et livraison offerte")
    add_footer(slide, page_num, total)

    # 3 scenario cards
    scenarios = renta.get("narrative_scenarios", [])
    colors = [RED, ORANGE, GREEN]
    for i, sc in enumerate(scenarios[:3]):
        x = Inches(0.4 + i * 4.2)
        y = Inches(1.4)
        w = Inches(3.9)
        h = Inches(2.2)

        border_c = colors[i]
        card = add_shape(slide, x, y, w, h, fill_color=WHITE, border_color=border_c, border_width=Pt(2))
        add_shape(slide, x, y, w, Inches(0.06), fill_color=border_c)

        add_textbox(slide, x + Inches(0.15), y + Inches(0.12), w - Inches(0.3), Inches(0.35),
                    sc.get("label", f"Scénario {i+1}"), font_size=12, color=MARINE, bold=True, alignment=PP_ALIGN.CENTER)

        contrib = sc.get("contribution_nette", 0)
        contrib_color = GREEN if contrib > 5 else (ORANGE if contrib > 0 else RED)

        add_textbox(slide, x + Inches(0.15), y + Inches(0.5), w - Inches(0.3), Inches(0.45),
                    f"{contrib:.2f}€", font_size=28, color=contrib_color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.95), w - Inches(0.3), Inches(0.25),
                    "Contribution nette / commande", font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

        details = [
            f"Prix : {sc['price']}€ → {sc['net_price']}€ (-{sc['discount_pct']}%)",
            f"Marge brute : {sc['gross_margin']:.2f}€ ({sc['margin_rate_on_original']}%)",
            f"Logistique : -{sc['logistics_cost']}€",
        ]
        for j, d in enumerate(details):
            add_textbox(slide, x + Inches(0.15), y + Inches(1.25 + j * 0.25), w - Inches(0.3), Inches(0.25),
                        d, font_size=9, color=DARK_GRAY)

    # Full matrix summary
    rows_matrix = [["Prix panier", "Marge 40%\nLogistique 4€", "Marge 45%\nLogistique 5€", "Marge 50%\nLogistique 7€"]]
    for item in renta.get("full_matrix", []):
        price = item["price"]
        m40 = item.get("m40_l4", {}).get("contribution_nette", "N/A")
        m45 = item.get("m45_l5", {}).get("contribution_nette", "N/A")
        m50 = item.get("m50_l7", {}).get("contribution_nette", "N/A")

        def fmt(v):
            if isinstance(v, (int, float)):
                return f"{v:.2f}€" + (" ⚠️" if v < 2 else (" ✅" if v > 5 else ""))
            return str(v)

        rows_matrix.append([f"{price}€", fmt(m40), fmt(m45), fmt(m50)])

    col_widths = [Inches(1.5), Inches(2.5), Inches(2.5), Inches(2.5)]
    build_table(slide, Inches(0.4), Inches(3.9), Inches(9), rows_matrix, col_widths)

    # Conclusion box
    threshold = renta.get("recommended_threshold_rounded", 25)
    add_shape(slide, Inches(9.7), Inches(3.9), Inches(3.4), Inches(1.6),
              fill_color=LIGHT_BLUE, border_color=MARINE, border_width=Pt(2))
    add_textbox(slide, Inches(9.85), Inches(4.0), Inches(3.1), Inches(1.4), "")
    tb = slide.shapes[-1]
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Seuil recommandé"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MARINE
    add_paragraph(tf, f"{threshold}€ minimum", font_size=24, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER, space_before=Pt(6))
    add_paragraph(tf, "pour activer l'abonnement\navec livraison offerte", font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(4))


def slide_reco_products(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_header_bar(slide, "Recommandation Impulse — Produits éligibles", "Sélection basée sur la récurrence naturelle d'achat")
    add_footer(slide, page_num, total)

    # 3 columns: ✅ Go / ⚠️ Attention / ❌ Exclure
    categories = [
        ("✅ Candidats naturels", GREEN,
         "Produits à consommation régulière et prévisible, faible complexité SKU.",
         ["Créatine (Creapure) — consommation quotidienne, 1 seul SKU",
          "Oméga-3 — cure continue, format capsules universel",
          "Magnésium Bisglycinate — cure longue durée",
          "Collagène (neutre / sans goût) — pas de problème de goût",
          "Multivitamines — consommation quotidienne",
          "Vitamine D3 — cure hivernale/continue",
          "Électrolytes — coureurs réguliers"]),
        ("⚠️ À traiter avec précaution", ORANGE,
         "Produits avec variantes qui nécessitent personnalisation.",
         ["Whey / Protéines — risque lassitude goût, besoin de variété",
          "Pre-workout — consommation irrégulière, dépend du rythme",
          "Barres protéinées — risque goût + stock important",
          "→ Solution : permettre le changement de goût à chaque cycle",
          "→ Ou proposer uniquement les formats neutres en abo"]),
        ("❌ À exclure au lancement", RED,
         "Produits saisonniers ou à renouvellement fréquent.",
         ["Produits saisonniers (gel été, boisson été)",
          "Formats découverte / échantillons",
          "Accessoires (shakers, gourdes)",
          "Produits à formule en cours de changement"])
    ]

    for i, (title, color, desc, items) in enumerate(categories):
        x = Inches(0.4 + i * 4.2)
        y = Inches(1.4)
        w = Inches(3.9)

        # Header
        add_shape(slide, x, y, w, Inches(0.45), fill_color=color)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.35),
                    title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, x + Inches(0.1), y + Inches(0.55), w - Inches(0.2), Inches(0.5),
                    desc, font_size=9, color=MID_GRAY)

        # Items
        tb = add_textbox(slide, x + Inches(0.1), y + Inches(1.05), w - Inches(0.2), Inches(4), "")
        tf = tb.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            if j == 0:
                tf.paragraphs[0].text = f"• {item}"
                tf.paragraphs[0].font.size = Pt(10)
                tf.paragraphs[0].font.color.rgb = DARK_GRAY
            else:
                add_paragraph(tf, f"• {item}", font_size=10, color=DARK_GRAY, space_before=Pt(5))

    # Bottom insight
    add_shape(slide, Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.9),
              fill_color=LIGHT_BLUE, border_color=MARINE, border_width=Pt(1))
    add_textbox(slide, Inches(0.7), Inches(6.05), Inches(12), Inches(0.8),
                "💡 Benchmark : Nutrimuscle propose ~250 produits en abo (quasi tout le catalogue). Aqeelab = 23/23 produits. "
                "Novoma = tout le catalogue. → Le marché tend vers l'exhaustivité. Impulse pourrait démarrer avec "
                "les 5-7 produits à haute récurrence puis élargir progressivement.",
                font_size=10, color=MARINE)


def slide_reco_model(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_header_bar(slide, "Recommandation Impulse — Modèle commercial", "Remise · Fréquences · Livraison · Engagement")
    add_footer(slide, page_num, total)

    # 4 pillar cards
    pillars = [
        ("Remise", "-10%", "Sur tous les produits en abonnement",
         ["Médiane marché : 10-15%", "Cohérent avec marge Impulse (45%)",
          "Moins agressif que Aqeelab (20%) mais soutenable",
          "Alternative : -10% 1ère, -15% dès la 2ème (modèle Nutrimuscle)"]),
        ("Fréquences", "1 · 2 · 3 mois", "3 fréquences au choix",
         ["Standard du marché (6/8 marques offrent 1-3 mois)",
          "1 mois = consommation quotidienne (créatine, oméga-3)",
          "2 mois = consommation modérée",
          "3 mois = stock (option économique)"]),
        ("Livraison", "Gratuite dès 25€", "En point relais pour les abonnés",
         ["Coût relais estimé : 4-5€ HT",
          "Seuil 25€ = rentable avec marge 45%",
          "Domicile : gratuit dès 49€ (vs 69€ standard)",
          "Différenciation comme Novoma (seuil réduit abo)"]),
        ("Engagement", "Sans engagement", "Annulation à tout moment",
         ["Standard absolu : 8/8 marques = sans engagement",
          "Pause & modification incluses",
          "Flexibilité ↓ churn de 15-18% (data marché)",
          "Réengagement par email si annulation"]),
    ]

    for i, (title, value, subtitle, details) in enumerate(pillars):
        x = Inches(0.3 + i * 3.25)
        y = Inches(1.4)
        w = Inches(3.05)

        card = add_shape(slide, x, y, w, Inches(4.5),
                         fill_color=WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD), border_width=Pt(1))
        add_shape(slide, x, y, w, Inches(0.06), fill_color=ORANGE)

        add_textbox(slide, x + Inches(0.1), y + Inches(0.12), w - Inches(0.2), Inches(0.3),
                    title, font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.4), w - Inches(0.2), Inches(0.5),
                    value, font_size=22, color=MARINE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.9), w - Inches(0.2), Inches(0.35),
                    subtitle, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

        # Divider
        add_shape(slide, x + Inches(0.3), y + Inches(1.3), w - Inches(0.6), Inches(0.015), fill_color=LIGHT_GRAY)

        tb = add_textbox(slide, x + Inches(0.15), y + Inches(1.4), w - Inches(0.3), Inches(3), "")
        tf = tb.text_frame
        tf.word_wrap = True
        for j, d in enumerate(details):
            if j == 0:
                tf.paragraphs[0].text = f"• {d}"
                tf.paragraphs[0].font.size = Pt(9)
                tf.paragraphs[0].font.color.rgb = DARK_GRAY
            else:
                add_paragraph(tf, f"• {d}", font_size=9, color=DARK_GRAY, space_before=Pt(4))

    # Bottom summary
    add_shape(slide, Inches(0.3), Inches(6.1), Inches(12.7), Inches(0.85),
              fill_color=MARINE)
    add_textbox(slide, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.75),
                "Résumé : Abonnement -10% · Sans engagement · 3 fréquences (1/2/3 mois) · Livraison gratuite dès 25€ en relais · "
                "Pause & modification possibles · Démarrage sur 5-7 produits à haute récurrence",
                font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


def slide_reco_ux(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_header_bar(slide, "Recommandation Impulse — UX & Positionnement", "Wording · Placement · Mise en avant")
    add_footer(slide, page_num, total)

    # Wording suggestion
    add_shape(slide, Inches(0.4), Inches(1.4), Inches(6), Inches(1.5),
              fill_color=LIGHT_BLUE, border_color=MARINE, border_width=Pt(2))
    add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.6), Inches(0.3),
                "Wording suggéré", font_size=14, color=MARINE, bold=True)
    add_textbox(slide, Inches(0.6), Inches(1.85), Inches(5.6), Inches(0.9),
                "CTA principal : \"S'abonner & économiser -10%\"\n"
                "Pitch fiche produit : \"Recevez automatiquement · Économisez 10% · Sans engagement\"\n"
                "Page dédiée : \"L'abonnement Impulse — Vos essentiels livrés, sans y penser\"",
                font_size=11, color=DARK_GRAY)

    # Placement module
    add_shape(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(1.5),
              fill_color=WHITE, border_color=RGBColor(0xDD, 0xDD, 0xDD), border_width=Pt(1))
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.6), Inches(0.3),
                "Placement sur la fiche produit", font_size=14, color=MARINE, bold=True)
    add_textbox(slide, Inches(7), Inches(1.85), Inches(5.6), Inches(0.9),
                "Toggle \"Achat unique / S'abonner\" au-dessus du CTA\n"
                "Prix barré + prix abonné + montant économisé en € (cf. Aqeelab)\n"
                "Badge \"Sans engagement\" visible (cf. Novoma)",
                font_size=11, color=DARK_GRAY)

    # Action plan grid
    actions = [
        ("Fiche produit", "Toggle abo au-dessus du CTA, prix barré, badge -10%", "P0 — Lancement", ORANGE),
        ("Page dédiée /abonnement", "3 étapes illustrées + avantages + FAQ (cf. Novoma)", "P0 — Lancement", ORANGE),
        ("Homepage", "Bannière ou section \"Abonnez-vous & économisez\"", "P1 — Mois 1", GREEN),
        ("Navigation", "Lien \"Abonnement\" dans le menu principal", "P1 — Mois 1", GREEN),
        ("Email post-achat", "\"Passez à l'abonnement et économisez sur votre prochain achat\"", "P1 — Mois 1", GREEN),
        ("Espace abonné", "Page Mon Compte avec gestion abo (cf. Nutrimuscle)", "P2 — Mois 2-3", MARINE),
        ("Email de rétention", "Email J-3 avant renouvellement + email anti-churn si annulation", "P2 — Mois 2-3", MARINE),
    ]

    rows = [["Touchpoint", "Action", "Priorité"]]
    for a in actions:
        rows.append([a[0], a[1], a[2]])

    col_widths = [Inches(2.5), Inches(6.5), Inches(2.5)]
    build_table(slide, Inches(0.4), Inches(3.2), Inches(12.5), rows, col_widths)

    # Bottom note
    add_textbox(slide, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.4),
                "💡 S'inspirer de Aqeelab pour l'affichage du montant économisé (€) et de Novoma pour l'argumentation santé sur la page dédiée.",
                font_size=10, color=MARINE)


def slide_next_steps(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_header_bar(slide, "Prochaines étapes", "Roadmap de mise en œuvre")
    add_footer(slide, page_num, total)

    steps = [
        ("Phase 1 : Validation", "Semaine 1-2",
         ["Valider le modèle commercial (remise, fréquences, seuils) avec le comité",
          "Confirmer les produits éligibles au lancement (5-7 SKUs)",
          "Valider le seuil panier minimum (25€)",
          "Arbitrage livraison : gratuite dès 25€ ou seuil plus élevé ?"]),
        ("Phase 2 : Setup technique", "Semaine 3-6",
         ["Choisir l'app Shopify : Recharge vs Shopify Native Subscriptions vs Bold",
          "Configurer le module abonnement sur les fiches produit",
          "Créer la page dédiée /abonnement",
          "Paramétrer les emails transactionnels (confirmation, rappel, renouvellement)"]),
        ("Phase 3 : Lancement soft", "Semaine 7-8",
         ["Lancement sur 5-7 produits pilotes",
          "Test A/B : avec vs sans abonnement sur les fiches produit",
          "Monitoring : taux de conversion abo, panier moyen, churn M1",
          "Ajustements UX/wording basés sur les premiers retours"]),
        ("Phase 4 : Scale", "Mois 3+",
         ["Extension à l'ensemble du catalogue éligible",
          "Mise en avant homepage + emails + navigation",
          "Programme de rétention (email anti-churn, crédit fidélité)",
          "Analyse LTV abonnés vs one-shot → ajustement remise si nécessaire"]),
    ]

    for i, (title, timing, items) in enumerate(steps):
        x = Inches(0.3 + i * 3.25)
        y = Inches(1.4)
        w = Inches(3.05)

        # Numbered circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.15), y, Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ORANGE
        circle.line.fill.background()
        tf_c = circle.text_frame
        tf_c.paragraphs[0].text = str(i + 1)
        tf_c.paragraphs[0].font.size = Pt(18)
        tf_c.paragraphs[0].font.color.rgb = WHITE
        tf_c.paragraphs[0].font.bold = True
        tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide, x, y + Inches(0.6), w, Inches(0.35),
                    title, font_size=13, color=MARINE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.95), w, Inches(0.25),
                    timing, font_size=10, color=ORANGE, alignment=PP_ALIGN.CENTER)

        # Divider
        add_shape(slide, x + Inches(0.3), y + Inches(1.25), w - Inches(0.6), Inches(0.015), fill_color=LIGHT_GRAY)

        tb = add_textbox(slide, x + Inches(0.1), y + Inches(1.35), w - Inches(0.2), Inches(3.5), "")
        tf = tb.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            if j == 0:
                tf.paragraphs[0].text = f"• {item}"
                tf.paragraphs[0].font.size = Pt(9)
                tf.paragraphs[0].font.color.rgb = DARK_GRAY
            else:
                add_paragraph(tf, f"• {item}", font_size=9, color=DARK_GRAY, space_before=Pt(5))


def slide_appendix_trustpilot(prs, master, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_header_bar(slide, "Annexe — Avis clients & Trustpilot", "Scores · Volumes · Forces & faiblesses perçues")
    add_footer(slide, page_num, total)

    brand_order = ["aqeelab", "nutripure", "novoma", "cuure", "nutriandco", "nutrimuscle", "myprotein", "decathlon"]

    rows = [["Marque", "Score", "Nb avis", "Points forts", "Points faibles"]]
    for bk in brand_order:
        bd = master["brands"].get(bk, {})
        avis = bd.get("avis", {})
        cat = bd.get("catalogue", {})
        name = cat.get("name", bk)
        score = f"{avis.get('trustpilot_score', '?')}/5"
        nb = f"{avis.get('trustpilot_nb_avis', '?'):,}" if isinstance(avis.get("trustpilot_nb_avis"), int) else "?"
        forts = ", ".join(avis.get("points_forts", [])[:2])
        faibles = ", ".join(avis.get("points_faibles", [])[:2])
        rows.append([name, score, nb, forts, faibles])

    col_widths = [Inches(1.5), Inches(0.8), Inches(1), Inches(4.5), Inches(4.5)]
    build_table(slide, Inches(0.4), Inches(1.4), Inches(12.5), rows, col_widths)

    add_textbox(slide, Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.5),
                "💡 Les marques premium françaises (Aqeelab, Nutripure, Novoma) dominent le classement Trustpilot. "
                "Les acteurs volume (MyProtein, Decathlon) souffrent sur la perception qualité/SAV.",
                font_size=10, color=MARINE)


def main():
    print("Loading data...")
    master, renta = load_data()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL_SLIDES = 20

    print("Building slides...")

    # 1. Title
    slide_title(prs, TOTAL_SLIDES)
    print("  ✅ 1/20 — Titre")

    # 2. Context
    slide_context(prs, TOTAL_SLIDES)
    print("  ✅ 2/20 — Contexte & Objectif")

    # 3. Methodology
    slide_methodology(prs, TOTAL_SLIDES)
    print("  ✅ 3/20 — Méthodologie")

    # 4-11. Brand slides
    brand_order = [
        ("nutriandco", 4),
        ("nutrimuscle", 5),
        ("novoma", 6),
        ("nutripure", 7),
        ("cuure", 8),
        ("decathlon", 9),
        ("aqeelab", 10),
        ("myprotein", 11),
    ]
    for bk, pg in brand_order:
        bd = master["brands"].get(bk, {})
        slide_brand(prs, bk, bd, pg, TOTAL_SLIDES)
        name = bd.get("catalogue", {}).get("name", bk)
        print(f"  ✅ {pg}/20 — {name}")

    # 12. Synthesis table
    slide_synthesis_table(prs, master, 12, TOTAL_SLIDES)
    print("  ✅ 12/20 — Synthèse comparative")

    # 13. UX comparative
    slide_ux_comparative(prs, master, 13, TOTAL_SLIDES)
    print("  ✅ 13/20 — UX comparative")

    # 14. Logistics
    slide_logistics(prs, master, 14, TOTAL_SLIDES)
    print("  ✅ 14/20 — Logistique")

    # 15. Profitability
    slide_profitability(prs, renta, 15, TOTAL_SLIDES)
    print("  ✅ 15/20 — Rentabilité")

    # 16. Reco products
    slide_reco_products(prs, 16, TOTAL_SLIDES)
    print("  ✅ 16/20 — Reco: Produits éligibles")

    # 17. Reco model
    slide_reco_model(prs, 17, TOTAL_SLIDES)
    print("  ✅ 17/20 — Reco: Modèle commercial")

    # 18. Reco UX
    slide_reco_ux(prs, 18, TOTAL_SLIDES)
    print("  ✅ 18/20 — Reco: UX & Positionnement")

    # 19. Next steps
    slide_next_steps(prs, 19, TOTAL_SLIDES)
    print("  ✅ 19/20 — Prochaines étapes")

    # 20. Appendix
    slide_appendix_trustpilot(prs, master, 20, TOTAL_SLIDES)
    print("  ✅ 20/20 — Annexe: Trustpilot")

    # Save
    output_path = os.path.join(ROOT, "Benchmark_Abonnement_Impulse_V2.pptx")
    prs.save(output_path)
    file_size = os.path.getsize(output_path)
    print(f"\n🎉 Deck V2 saved: {output_path}")
    print(f"   {TOTAL_SLIDES} slides · {file_size / 1024 / 1024:.1f} MB")


if __name__ == "__main__":
    main()
